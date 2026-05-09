"""Generate the VidIQ presentation deck as a .pptx that imports cleanly into
Canva or opens in PowerPoint / Keynote.

Slide blueprint: 12-presentation-outline.md (15 slides + 1 backup).
Brand: dark midnight bg (#0a0612), violet→fuchsia→cyan accents.
Typography: Plus Jakarta Sans (display), Inter (body), JetBrains Mono (mono).

Re-run with:
    ./venv/Scripts/python.exe marketing/build_deck.py
Output:
    marketing/submissions/VidIQ_Final_Presentation.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# ── Brand palette (matches frontend/src/app/globals.css) ───────────────────
BG_DEEP = RGBColor(0x0A, 0x06, 0x12)      # midnight
BG_CARD = RGBColor(0x16, 0x10, 0x21)      # raised card on bg
FG_HIGH = RGBColor(0xF5, 0xEF, 0xFF)      # primary text on dark
FG_MID = RGBColor(0xB8, 0xB0, 0xC8)       # muted text on dark
FG_LOW = RGBColor(0x6B, 0x60, 0x80)       # tertiary text
ACCENT_VIOLET = RGBColor(0xA8, 0x55, 0xF7)
ACCENT_FUCHSIA = RGBColor(0xEC, 0x48, 0x99)
ACCENT_CYAN = RGBColor(0x06, 0xB6, 0xD4)
ACCENT_EMERALD = RGBColor(0x10, 0xB9, 0x81)
ACCENT_AMBER = RGBColor(0xF5, 0x9E, 0x0B)

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Type stack — Canva will substitute if these aren't installed
DISPLAY = "Plus Jakarta Sans"
BODY = "Inter"
MONO = "JetBrains Mono"


# ── Asset paths ────────────────────────────────────────────────────────────

ASSETS = Path(__file__).resolve().parent / "assets"
SHOTS = ASSETS / "screenshots"
FIREFLY = ASSETS / "firefly"
VIDEO = ASSETS / "video-ad"
SEO = ASSETS / "seo-evidence"
# Brand logos are sourced from the frontend's public dir so the site +
# the deck stay in lockstep (single source of truth).
PUBLIC = Path(__file__).resolve().parent.parent / "frontend" / "public"


def asset(*names: str) -> Path | None:
    """Return the first path that exists for any of the given names."""
    for n in names:
        for root in (SHOTS, FIREFLY, VIDEO, SEO, PUBLIC):
            p = root / n
            if p.exists():
                return p
    return None


# Convenience aliases for the 3 brand-logo lockups
LOGO_MARK_DARK = "vidiq_logo_black_bg.png"   # square mark on dark bg (use on dark slides)
LOGO_MARK_LIGHT = "vidiq_logo_white_bg.png"  # square mark on light bg
LOGO_TEXT = "vidiq_logo_text.png"            # horizontal wordmark lockup


# ── Helpers ────────────────────────────────────────────────────────────────


def add_image(slide, src: Path, *, left, top, width=None, height=None,
              border_color=None, border_pt=1.0):
    """Place an image on the slide. Adds a thin tinted frame so screenshots
    don't float in space against the dark background."""
    if not src or not src.exists():
        # Drop a placeholder card if asset missing
        ph = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
            width or Inches(2), height or Inches(2),
        )
        ph.fill.solid()
        ph.fill.fore_color.rgb = BG_CARD
        ph.line.color.rgb = ACCENT_FUCHSIA
        ph.line.width = Pt(0.75)
        ph.adjustments[0] = 0.05
        tf = ph.text_frame
        tf.margin_left = Inches(0.15)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f"missing\n{src.name if src else '?'}"
        r.font.name = MONO
        r.font.size = Pt(9)
        r.font.color.rgb = ACCENT_FUCHSIA
        return ph

    pic = slide.shapes.add_picture(
        str(src), left, top,
        width=width, height=height,
    )
    if border_color is not None:
        pic.line.color.rgb = border_color
        pic.line.width = Pt(border_pt)
    return pic


def add_blank_slide(prs: Presentation):
    blank = prs.slide_layouts[6]  # 6 is the blank layout in default theme
    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_DEEP
    bg.shadow.inherit = False
    return slide


def add_text(slide, text, *, left, top, width, height,
             size=18, bold=False, font=BODY, color=FG_HIGH,
             align=PP_ALIGN.LEFT, line_spacing=1.2):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tf


def add_bullets(slide, bullets, *, left, top, width, height,
                size=16, font=BODY, color=FG_HIGH, accent=ACCENT_VIOLET):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.4
        p.space_after = Pt(8)
        # Coloured square bullet
        sq = p.add_run()
        sq.text = "■  "
        sq.font.name = font
        sq.font.size = Pt(size)
        sq.font.color.rgb = accent
        run = p.add_run()
        run.text = b
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tf


def add_accent_bar(slide, *, top, height=Inches(0.06), color=ACCENT_VIOLET, left=Inches(0.7), width=Inches(2.5)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = color


def add_header(slide, eyebrow, title, *, accent=ACCENT_VIOLET):
    add_text(slide, eyebrow.upper(),
             left=Inches(0.7), top=Inches(0.6), width=Inches(8), height=Inches(0.3),
             size=10, bold=True, color=accent, font=BODY)
    add_text(slide, title,
             left=Inches(0.7), top=Inches(0.9), width=Inches(12), height=Inches(1.2),
             size=40, bold=True, font=DISPLAY, color=FG_HIGH, line_spacing=1.05)
    add_accent_bar(slide, top=Inches(2.05), color=accent)


def add_footer(slide, pillar):
    # Small brand-mark in the bottom-left corner — same square logo used in the
    # site favicon, so the deck visually aligns with the live product.
    mark = asset(LOGO_MARK_DARK)
    if mark:
        add_image(slide, mark,
                  left=Inches(0.7), top=Inches(7.0),
                  width=Inches(0.32), height=Inches(0.32))
        text_left = Inches(1.1)
    else:
        text_left = Inches(0.7)
    add_text(slide, "VidIQ",
             left=text_left, top=Inches(7.05), width=Inches(1.6), height=Inches(0.3),
             size=9, bold=True, font=DISPLAY, color=ACCENT_VIOLET)
    add_text(slide, pillar,
             left=Inches(2.7), top=Inches(7.05), width=Inches(10), height=Inches(0.3),
             size=9, font=BODY, color=FG_LOW)


def add_card(slide, *, left, top, width, height, fill=BG_CARD, accent=None):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    c.fill.solid()
    c.fill.fore_color.rgb = fill
    c.line.color.rgb = accent or RGBColor(0x2A, 0x22, 0x3A)
    c.line.width = Pt(0.75)
    c.adjustments[0] = 0.06
    return c


def add_chip(slide, text, *, left, top, width=Inches(1.5), height=Inches(0.32),
             accent=ACCENT_VIOLET, fill=None):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    c.fill.solid()
    c.fill.fore_color.rgb = fill or RGBColor(0x1F, 0x12, 0x3A)
    c.line.color.rgb = accent
    c.line.width = Pt(0.5)
    c.adjustments[0] = 0.5
    tf = c.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = BODY
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = accent


def set_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


# ── Slide builders ─────────────────────────────────────────────────────────


def slide_1_cover(prs):
    s = add_blank_slide(prs)

    # ── Layered aurora glows (top-left + bottom-right) ────────────────
    glow_a = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-3.5), Inches(-3.5),
                                Inches(8), Inches(8))
    glow_a.fill.solid()
    glow_a.fill.fore_color.rgb = RGBColor(0x33, 0x12, 0x55)
    glow_a.line.fill.background()

    glow_b = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(4),
                                Inches(7), Inches(7))
    glow_b.fill.solid()
    glow_b.fill.fore_color.rgb = RGBColor(0x44, 0x18, 0x55)
    glow_b.line.fill.background()

    # Subtle grid stripes (top-left to bottom-right diagonal feel)
    for i, y_off in enumerate([0.5, 1.6, 2.7, 5.5, 6.4, 7.0]):
        stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(0), Inches(y_off),
                                    SLIDE_W, Inches(0.005))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = RGBColor(0x44, 0x33, 0x55)
        stripe.line.fill.background()

    # ── Brand mark watermark — top-left corner ───────────────────────
    mark = asset(LOGO_MARK_DARK)
    if mark:
        add_image(s, mark,
                  left=Inches(0.55), top=Inches(0.45),
                  width=Inches(0.5), height=Inches(0.5))

    # ── Eyebrow tag, top of slide ─────────────────────────────────────
    eyebrow_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(5.42), Inches(0.52),
                                      Inches(2.5), Inches(0.34))
    eyebrow_card.fill.solid()
    eyebrow_card.fill.fore_color.rgb = RGBColor(0x1F, 0x12, 0x3A)
    eyebrow_card.line.color.rgb = ACCENT_VIOLET
    eyebrow_card.line.width = Pt(0.5)
    eyebrow_card.adjustments[0] = 0.5
    add_text(s, "SEMESTER COURSE PROJECT",
             left=Inches(5.42), top=Inches(0.55), width=Inches(2.5), height=Inches(0.3),
             size=8, bold=True, font=BODY, color=ACCENT_VIOLET, align=PP_ALIGN.CENTER)

    # ── Wordmark ──────────────────────────────────────────────────────
    add_text(s, "VidIQ",
             left=Inches(0), top=Inches(1.05), width=SLIDE_W, height=Inches(1.6),
             size=128, bold=True, font=DISPLAY, color=FG_HIGH,
             align=PP_ALIGN.CENTER, line_spacing=1)

    # Triple-segment accent bar (violet · fuchsia · cyan) under wordmark
    seg_w = 0.8
    total_w = seg_w * 3 + 0.06 * 2
    start_x = (SLIDE_W.inches - total_w) / 2
    for i, c in enumerate([ACCENT_VIOLET, ACCENT_FUCHSIA, ACCENT_CYAN]):
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(start_x + i * (seg_w + 0.06)), Inches(2.78),
                                 Inches(seg_w), Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = c
        bar.line.fill.background()

    # Subtitle + tagline
    add_text(s, "AI Video Intelligence",
             left=Inches(0), top=Inches(2.95), width=SLIDE_W, height=Inches(0.5),
             size=22, bold=True, font=DISPLAY, color=ACCENT_VIOLET, align=PP_ALIGN.CENTER)
    add_text(s, "Watch less. Learn more.",
             left=Inches(0), top=Inches(3.45), width=SLIDE_W, height=Inches(0.5),
             size=16, font=BODY, color=FG_MID, align=PP_ALIGN.CENTER)

    # ── Team panel (left) + Course panel (right) ──────────────────────
    panel_y = Inches(4.25)
    panel_h = Inches(2.4)

    # Team — left card
    team_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(1.0), panel_y,
                                   Inches(6.0), panel_h)
    team_card.fill.solid()
    team_card.fill.fore_color.rgb = BG_CARD
    team_card.line.color.rgb = ACCENT_VIOLET
    team_card.line.width = Pt(0.75)
    team_card.adjustments[0] = 0.05

    add_text(s, "GROUP MEMBERS",
             left=Inches(1.3), top=panel_y + Inches(0.18),
             width=Inches(5.4), height=Inches(0.28),
             size=10, bold=True, font=BODY, color=ACCENT_VIOLET)

    members = [
        ("22i-1653", "Insharah Aman"),
        ("21i-0416", "M. Nouman Hafeez"),
        ("21i-0484", "Shayan Khan"),
        ("21i-2507", "Muhammad Zain"),
        ("22i-1200", "Farhan Ahmed"),
    ]
    for i, (rno, name) in enumerate(members):
        y = panel_y + Inches(0.55 + i * 0.36)
        # roll number badge
        add_text(s, rno,
                 left=Inches(1.3), top=y, width=Inches(1.4), height=Inches(0.32),
                 size=11, bold=True, font=MONO, color=ACCENT_CYAN)
        # name
        add_text(s, name,
                 left=Inches(2.8), top=y, width=Inches(3.9), height=Inches(0.32),
                 size=14, font=BODY, color=FG_HIGH)

    # Course — right card
    course_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(7.3), panel_y,
                                     Inches(5.0), panel_h)
    course_card.fill.solid()
    course_card.fill.fore_color.rgb = BG_CARD
    course_card.line.color.rgb = ACCENT_FUCHSIA
    course_card.line.width = Pt(0.75)
    course_card.adjustments[0] = 0.05

    add_text(s, "COURSE",
             left=Inches(7.6), top=panel_y + Inches(0.18),
             width=Inches(4.5), height=Inches(0.28),
             size=10, bold=True, font=BODY, color=ACCENT_FUCHSIA)

    course_rows = [
        ("Course",     "Digital Marketing"),
        ("Section",    "CS-A"),
        ("Instructor", "Sir Maaz Zafar Cheema"),
        ("Project",    "Semester Course Project"),
    ]
    for i, (label, value) in enumerate(course_rows):
        y = panel_y + Inches(0.55 + i * 0.42)
        add_text(s, label.upper(),
                 left=Inches(7.6), top=y, width=Inches(1.3), height=Inches(0.3),
                 size=9, bold=True, font=BODY, color=FG_LOW)
        add_text(s, value,
                 left=Inches(8.95), top=y - Inches(0.02),
                 width=Inches(3.3), height=Inches(0.34),
                 size=14, bold=(label == "Project"), font=DISPLAY if label == "Project" else BODY,
                 color=FG_HIGH if label == "Project" else FG_HIGH)

    # ── Bottom URL pill (centered) ────────────────────────────────────
    url_pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(4.42), Inches(7.0),
                                  Inches(4.5), Inches(0.4))
    url_pill.fill.solid()
    url_pill.fill.fore_color.rgb = RGBColor(0x1F, 0x12, 0x3A)
    url_pill.line.color.rgb = ACCENT_CYAN
    url_pill.line.width = Pt(0.6)
    url_pill.adjustments[0] = 0.5
    add_text(s, "🔗  vidiq-two.vercel.app",
             left=Inches(4.42), top=Inches(7.05), width=Inches(4.5), height=Inches(0.32),
             size=12, font=MONO, color=ACCENT_CYAN, align=PP_ALIGN.CENTER, bold=True)

    # ── Decorative floating orbs (small, on the diagonals) ────────────
    for cx, cy, dia, col in [
        (0.6, 0.6, 0.18, ACCENT_VIOLET),
        (12.6, 0.5, 0.12, ACCENT_FUCHSIA),
        (0.4, 6.95, 0.14, ACCENT_CYAN),
        (12.7, 7.05, 0.16, ACCENT_VIOLET),
    ]:
        orb = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(cx - dia / 2), Inches(cy - dia / 2),
                                 Inches(dia), Inches(dia))
        orb.fill.solid()
        orb.fill.fore_color.rgb = col
        orb.line.fill.background()

    # ── Apply a slide-level fade transition (XML insertion) ───────────
    try:
        from lxml import etree
        from pptx.oxml.ns import qn
        s_elem = s.element
        # PPTX transition XML
        trans_xml = (
            '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
            'spd="med">'
            '<p:fade/>'
            '</p:transition>'
        )
        trans = etree.fromstring(trans_xml)
        s_elem.append(trans)
    except Exception:
        pass

    set_notes(s,
        "Open with the hook from the ad: 'You hit play on a 90-min video and walked away three "
        "minutes later.' Pause. Introduce the team in one breath, then click straight to slide 2. "
        "Time: 30 seconds.")
    return s


def slide_2_problem(prs):
    s = add_blank_slide(prs)
    add_header(s, "Problem", "Long-form video is the most-uploaded but least-read medium.", accent=ACCENT_FUCHSIA)
    add_footer(s, "Pillar 0 · Problem framing")

    bullets = [
        "500 hours of video uploaded to YouTube every minute (Statista, 2024).",
        "Average watch time per session: 17 minutes — a 90-minute lecture takes 5 sessions.",
        "65% of students rewatch lecture videos at 2× — search and citations are the bottleneck.",
        "Existing tools summarise text (NotebookLM) or audio (Otter); none combine speech + vision.",
    ]
    add_bullets(s, bullets,
                left=Inches(0.7), top=Inches(2.5), width=Inches(12), height=Inches(4),
                size=18, accent=ACCENT_FUCHSIA)

    set_notes(s, "Read all 4 stats out loud. Pause after the third — that's the strongest hook for the marker. "
                 "Time: 90 seconds.")


def slide_3_product(prs):
    s = add_blank_slide(prs)
    add_header(s, "Pillar 3 · Product", "The dashboard you'll see in the live demo.", accent=ACCENT_VIOLET)
    add_footer(s, "Pillar 3 · Product")

    # ── Left half: introductory walkthrough video (click to play) ────
    intro = asset("introductory-video.mp4")
    poster = asset(LOGO_MARK_DARK)
    video_left = Inches(0.7)
    video_top = Inches(2.4)
    video_w = Inches(6.5)
    video_h = Inches(3.65)  # 16:9 ratio of video_w
    add_card(s, left=video_left, top=video_top,
             width=video_w, height=video_h + Inches(0.55), accent=ACCENT_VIOLET)
    if intro:
        try:
            s.shapes.add_movie(
                str(intro),
                video_left + Inches(0.15), video_top + Inches(0.15),
                video_w - Inches(0.3), video_h,
                poster_frame_image=str(poster) if poster else None,
                mime_type="video/mp4",
            )
        except Exception:
            add_text(s, "▶  introductory-video.mp4 — open from PowerPoint to play",
                     left=video_left + Inches(0.15), top=video_top + Inches(1.5),
                     width=video_w - Inches(0.3), height=Inches(0.5),
                     size=12, font=MONO, color=ACCENT_VIOLET, align=PP_ALIGN.CENTER)
    else:
        add_text(s, "[ intro walkthrough video — drop into marketing/assets/video-ad/ ]",
                 left=video_left + Inches(0.15), top=video_top + Inches(1.6),
                 width=video_w - Inches(0.3), height=Inches(0.5),
                 size=11, font=MONO, color=ACCENT_FUCHSIA, align=PP_ALIGN.CENTER)
    # Caption strip
    add_text(s, "▶  Walkthrough · paste a YouTube URL, watch the pipeline run end-to-end",
             left=video_left + Inches(0.15), top=video_top + video_h + Inches(0.18),
             width=video_w - Inches(0.3), height=Inches(0.3),
             size=10, font=BODY, color=FG_MID, align=PP_ALIGN.CENTER)

    # ── Right half: 4 feature cards (compact 2x2 grid) ───────────────
    features = [
        ("Multimodal pipeline",  "Speech (Whisper) + vision (Gemini) + LLM, fused.",      ACCENT_VIOLET),
        ("Time-stamped citations","Every claim cites the exact moment. Click to seek.",    ACCENT_CYAN),
        ("Live + recorded",       "YouTube videos AND live streams — neither competitor.", ACCENT_FUCHSIA),
        ("Domain-aware",          "Medical / legal / trading / education prompt modes.",   ACCENT_EMERALD),
    ]
    fcard_w = Inches(2.85); fcard_h = Inches(1.92); fgap = Inches(0.12)
    fsx = Inches(7.45); fsy = Inches(2.4)
    for i, (title, body, accent) in enumerate(features):
        col = i % 2; row = i // 2
        x = fsx + (fcard_w + fgap) * col
        y = fsy + (fcard_h + fgap) * row
        add_card(s, left=x, top=y, width=fcard_w, height=fcard_h, accent=accent)
        add_text(s, title,
                 left=x + Inches(0.22), top=y + Inches(0.18),
                 width=fcard_w - Inches(0.4), height=Inches(0.55),
                 size=13, bold=True, font=DISPLAY, color=FG_HIGH, line_spacing=1.15)
        add_text(s, body,
                 left=x + Inches(0.22), top=y + Inches(0.78),
                 width=fcard_w - Inches(0.4), height=Inches(1.1),
                 size=10, font=BODY, color=FG_MID, line_spacing=1.3)

    add_text(s, "LIVE  →  vidiq-two.vercel.app",
             left=Inches(0.7), top=Inches(6.7), width=Inches(12), height=Inches(0.35),
             size=12, font=MONO, color=ACCENT_CYAN, bold=True)

    set_notes(s, "Click the video on the left to play the walkthrough. If the live demo stalls, "
                 "the recorded walkthrough is your fallback. The 4 feature cards on the right are "
                 "the talking points to read while the video plays. Time: 75 seconds.")


def slide_4_pipeline(prs):
    s = add_blank_slide(prs)
    add_header(s, "Pillar 3 · Architecture", "End-to-end multimodal pipeline.", accent=ACCENT_VIOLET)
    add_footer(s, "Pillar 3 · Architecture")

    # Pipeline boxes left-to-right
    steps = [
        ("URL input", "yt-dlp resolve", ACCENT_VIOLET),
        ("Audio", "ffmpeg → Whisper", ACCENT_CYAN),
        ("Frames", "OpenCV scene-cut", ACCENT_EMERALD),
        ("Vision", "Gemini caption", ACCENT_AMBER),
        ("Synthesise", "LLM JSON", ACCENT_FUCHSIA),
    ]
    box_w = Inches(2.2)
    box_h = Inches(1.3)
    gap = Inches(0.25)
    total_w = box_w * len(steps) + gap * (len(steps) - 1)
    start_x = (SLIDE_W - total_w) / 2
    y = Inches(3.2)

    for i, (title, sub, accent) in enumerate(steps):
        x = start_x + (box_w + gap) * i
        add_card(s, left=x, top=y, width=box_w, height=box_h, accent=accent)
        add_text(s, title,
                 left=x + Inches(0.15), top=y + Inches(0.18), width=box_w - Inches(0.3), height=Inches(0.4),
                 size=14, bold=True, font=DISPLAY, color=FG_HIGH, align=PP_ALIGN.CENTER)
        add_text(s, sub,
                 left=x + Inches(0.15), top=y + Inches(0.7), width=box_w - Inches(0.3), height=Inches(0.5),
                 size=10, font=MONO, color=accent, align=PP_ALIGN.CENTER)
        # Arrow
        if i < len(steps) - 1:
            arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                       x + box_w + Inches(0.02), y + Inches(0.55),
                                       gap - Inches(0.04), Inches(0.2))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT_VIOLET
            arrow.line.fill.background()

    add_text(s, "Stack:  FastAPI · Next.js 14 · Gemini · faster-whisper · OpenCV · SQLite · WebSocket",
             left=Inches(0.7), top=Inches(5.5), width=Inches(12), height=Inches(0.4),
             size=12, font=MONO, color=FG_MID, align=PP_ALIGN.CENTER)

    add_text(s, "Provider-agnostic: Gemini ↔ Groq ↔ OpenAI auto-failover.  Free-tier capped — $0 to operate.",
             left=Inches(0.7), top=Inches(6), width=Inches(12), height=Inches(0.4),
             size=11, font=BODY, color=FG_LOW, align=PP_ALIGN.CENTER)

    set_notes(s, "Speak to the multimodal angle: 'transcript alone is text — frames + transcript is understanding.' "
                 "Mention the failover (real engineering moat). Time: 60 seconds.")


def slide_5_brand(prs):
    s = add_blank_slide(prs)
    add_header(s, "Pillar 1 · Branding", "Identity that signals AI without the cliché.", accent=ACCENT_FUCHSIA)
    add_footer(s, "Pillar 1 · Branding")

    # Two columns: left = palette swatches, right = type sample + voice
    # ── Palette
    add_text(s, "PALETTE",
             left=Inches(0.7), top=Inches(2.4), width=Inches(5), height=Inches(0.3),
             size=10, bold=True, font=BODY, color=FG_LOW)
    swatches = [
        ("Violet 500", ACCENT_VIOLET, "Trust + creativity"),
        ("Fuchsia 500", ACCENT_FUCHSIA, "Energy + signal"),
        ("Cyan 500", ACCENT_CYAN, "Clarity + intelligence"),
        ("Midnight", BG_DEEP, "Stage / canvas"),
    ]
    for i, (name, col, why) in enumerate(swatches):
        y = Inches(2.85 + i * 0.85)
        sq = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.7), y, Inches(0.7), Inches(0.7))
        sq.fill.solid()
        sq.fill.fore_color.rgb = col
        sq.line.color.rgb = RGBColor(0x44, 0x33, 0x55)
        sq.line.width = Pt(0.5)
        add_text(s, name,
                 left=Inches(1.6), top=y + Inches(0.05), width=Inches(2.5), height=Inches(0.35),
                 size=14, bold=True, font=DISPLAY, color=FG_HIGH)
        add_text(s, why,
                 left=Inches(1.6), top=y + Inches(0.4), width=Inches(2.5), height=Inches(0.3),
                 size=10, font=BODY, color=FG_MID)

    # ── Typography
    add_text(s, "TYPOGRAPHY",
             left=Inches(7), top=Inches(2.4), width=Inches(5), height=Inches(0.3),
             size=10, bold=True, font=BODY, color=FG_LOW)
    add_text(s, "Plus Jakarta Sans 800",
             left=Inches(7), top=Inches(2.8), width=Inches(6), height=Inches(0.7),
             size=32, bold=True, font=DISPLAY, color=FG_HIGH)
    add_text(s, "Display headings · humanist geometric",
             left=Inches(7), top=Inches(3.45), width=Inches(6), height=Inches(0.3),
             size=11, font=BODY, color=FG_LOW)
    add_text(s, "Inter 400 / 500 / 600",
             left=Inches(7), top=Inches(4), width=Inches(6), height=Inches(0.5),
             size=20, font=BODY, color=FG_HIGH)
    add_text(s, "Body copy · interface · long-form",
             left=Inches(7), top=Inches(4.5), width=Inches(6), height=Inches(0.3),
             size=11, font=BODY, color=FG_LOW)
    add_text(s, "JetBrains Mono",
             left=Inches(7), top=Inches(5.05), width=Inches(6), height=Inches(0.4),
             size=14, font=MONO, color=ACCENT_CYAN)
    add_text(s, "Code · timestamps · technical labels",
             left=Inches(7), top=Inches(5.5), width=Inches(6), height=Inches(0.3),
             size=11, font=BODY, color=FG_LOW)

    # Brand-consistency footnote — pointer to dedicated logo slide
    add_text(s, "Brand applied across 7 surfaces:  favicon · top-nav · hero · OG card · JSON-LD · meta · FB cover.",
             left=Inches(0.7), top=Inches(6.4), width=Inches(12), height=Inches(0.4),
             size=11, font=BODY, color=FG_MID, line_spacing=1.4, align=PP_ALIGN.CENTER)
    add_text(s, "Logo system →  next slide.",
             left=Inches(0.7), top=Inches(6.85), width=Inches(12), height=Inches(0.3),
             size=10, font=BODY, color=ACCENT_VIOLET, align=PP_ALIGN.CENTER, bold=True)

    set_notes(s, "Why violet, not blue (default SaaS): blue is the consultancy default — violet retains the trust association "
                 "while signalling creativity and AI. Cite Mehta-Zhu (2009) on red/blue cognition. "
                 "The logo system is presented in detail on the next slide. Time: 45 seconds.")


def slide_6_ad(prs):
    s = add_blank_slide(prs)
    add_header(s, "Pillar 1 · Creative", "30-second teaser ad — story arc.", accent=ACCENT_FUCHSIA)
    add_footer(s, "Pillar 1 · Creative")

    # Story-arc strip
    arc = [
        ("0–3 s", "Hook", "Pause the lecture", ACCENT_VIOLET),
        ("3–10 s", "Problem", "90 min, 17 min watched", ACCENT_FUCHSIA),
        ("10–20 s", "Product", "Paste URL → grounded summary", ACCENT_CYAN),
        ("20–26 s", "Proof", "Click any claim → exact moment", ACCENT_EMERALD),
        ("26–30 s", "CTA", "Watch less. Learn more.", ACCENT_AMBER),
    ]
    box_w = Inches(2.4)
    box_h = Inches(2.2)
    gap = Inches(0.15)
    total_w = box_w * len(arc) + gap * (len(arc) - 1)
    start_x = (SLIDE_W - total_w) / 2
    y = Inches(2.7)
    for i, (ts, name, body, accent) in enumerate(arc):
        x = start_x + (box_w + gap) * i
        add_card(s, left=x, top=y, width=box_w, height=box_h, accent=accent)
        add_text(s, ts,
                 left=x + Inches(0.18), top=y + Inches(0.18), width=box_w - Inches(0.4), height=Inches(0.3),
                 size=10, font=MONO, color=accent, bold=True)
        add_text(s, name,
                 left=x + Inches(0.18), top=y + Inches(0.55), width=box_w - Inches(0.4), height=Inches(0.5),
                 size=18, bold=True, font=DISPLAY, color=FG_HIGH)
        add_text(s, body,
                 left=x + Inches(0.18), top=y + Inches(1.05), width=box_w - Inches(0.4), height=Inches(1),
                 size=12, font=BODY, color=FG_MID, line_spacing=1.3)

    # ── Embed the 30s video (auto-rendered as a movie shape) ─────────
    video_file = asset("video_ad.mp4")
    poster = asset("IMG-8.png")  # static frame shown before play
    video_left = Inches(4.65)
    video_top = Inches(5.1)
    video_w = Inches(4.0)
    video_h = Inches(2.25)  # 16:9 default poster

    if video_file:
        try:
            s.shapes.add_movie(
                str(video_file),
                video_left, video_top, video_w, video_h,
                poster_frame_image=str(poster) if poster else None,
                mime_type="video/mp4",
            )
        except Exception:
            # Fallback — link the video as a text reference if embedding fails
            add_text(s, "▶  video_ad.mp4 (in marketing/assets/video-ad/) — open from PowerPoint to play",
                     left=video_left, top=video_top, width=video_w, height=video_h,
                     size=12, font=MONO, color=ACCENT_VIOLET, align=PP_ALIGN.CENTER)
    else:
        add_text(s, "[ video missing — re-export and place in marketing/assets/video-ad/video_ad.mp4 ]",
                 left=video_left, top=video_top, width=video_w, height=video_h,
                 size=12, font=MONO, color=ACCENT_FUCHSIA, align=PP_ALIGN.CENTER)

    add_text(s, "Cutdowns:  30s master · 15s social · 6s bumper · 10s in-stream",
             left=Inches(0.7), top=Inches(6.95), width=Inches(12), height=Inches(0.4),
             size=11, font=BODY, color=FG_LOW, align=PP_ALIGN.CENTER)

    set_notes(s, "Click the video to play it during presentation. The story-arc strip above stays "
                 "as a visual reference for the structure. Time: 35 seconds (the ad runs).")


def slide_7_social(prs):
    s = add_blank_slide(prs)
    add_header(s, "Pillar 2 · Social", "14-day calendar · FB Page · IG Business.", accent=ACCENT_CYAN)
    add_footer(s, "Pillar 2 · Social calendar")

    # Stat row at top — KPI #5 + #6 evidence numbers
    stats = [("28", "posts queued"), ("4", "formats"), ("7×", "weekly cadence"), ("2", "channels")]
    for i, (n, lbl) in enumerate(stats):
        x = Inches(0.7 + i * 1.65)
        add_card(s, left=x, top=Inches(2.4), width=Inches(1.5), height=Inches(1.05), accent=ACCENT_CYAN)
        add_text(s, n,
                 left=x, top=Inches(2.5), width=Inches(1.5), height=Inches(0.55),
                 size=28, bold=True, font=DISPLAY, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
        add_text(s, lbl.upper(),
                 left=x, top=Inches(3.05), width=Inches(1.5), height=Inches(0.3),
                 size=9, font=BODY, color=FG_MID, align=PP_ALIGN.CENTER)

    # Format chips on right of stat row
    add_text(s, "FORMATS",
             left=Inches(7.3), top=Inches(2.4), width=Inches(5), height=Inches(0.3),
             size=10, bold=True, font=BODY, color=FG_LOW)
    formats = ["Reels", "Carousels", "Single image", "Stories"]
    for i, f in enumerate(formats):
        add_chip(s, f, left=Inches(7.3 + (i % 2) * 1.55), top=Inches(2.7 + (i // 2) * 0.4),
                 width=Inches(1.45), accent=ACCENT_CYAN)

    # Three real screenshots: Planner left (large), FB page top-right, IG welcome bottom-right
    planner = asset("planner-calendar.png")
    fb_page = asset("fb-page.png")
    fb_welcome = asset("fb-welcome-post.png")
    ig_welcome = asset("insta-welcome-post.jpeg", "insta-welcome-post.png")

    # Planner — large left panel
    add_text(s, "BUSINESS SUITE PLANNER · 28 POSTS QUEUED",
             left=Inches(0.7), top=Inches(3.65), width=Inches(6.5), height=Inches(0.25),
             size=9, bold=True, font=BODY, color=ACCENT_CYAN)
    add_image(s, planner,
              left=Inches(0.7), top=Inches(3.95), width=Inches(6.5), height=Inches(2.95),
              border_color=ACCENT_CYAN, border_pt=1)

    # FB page — top right
    add_text(s, "FB PAGE · facebook.com/profile.php?id=61588939576479",
             left=Inches(7.3), top=Inches(3.65), width=Inches(2.6), height=Inches(0.25),
             size=8, bold=True, font=BODY, color=ACCENT_VIOLET)
    add_image(s, fb_page,
              left=Inches(7.3), top=Inches(3.95), width=Inches(2.6), height=Inches(1.4),
              border_color=ACCENT_VIOLET, border_pt=1)

    # FB welcome post + IG welcome post stacked
    add_text(s, "FB PINNED · welcome post",
             left=Inches(7.3), top=Inches(5.45), width=Inches(2.6), height=Inches(0.25),
             size=9, bold=True, font=BODY, color=ACCENT_FUCHSIA)
    add_image(s, fb_welcome,
              left=Inches(7.3), top=Inches(5.7), width=Inches(2.6), height=Inches(1.2),
              border_color=ACCENT_FUCHSIA, border_pt=1)

    add_text(s, "IG · @vidiq_official",
             left=Inches(10.1), top=Inches(3.65), width=Inches(2.55), height=Inches(0.25),
             size=9, bold=True, font=BODY, color=ACCENT_FUCHSIA)
    add_image(s, ig_welcome,
              left=Inches(10.1), top=Inches(3.95), width=Inches(2.55), height=Inches(2.95),
              border_color=ACCENT_FUCHSIA, border_pt=1)

    # Live profile URL pills (centered footer row)
    add_chip(s, "📘  fb.com/profile.php?id=61588939576479",
             left=Inches(2.7), top=Inches(6.95),
             width=Inches(3.6), height=Inches(0.32), accent=ACCENT_VIOLET)
    add_chip(s, "📸  instagram.com/vidiq_official/",
             left=Inches(6.7), top=Inches(6.95),
             width=Inches(3.6), height=Inches(0.32), accent=ACCENT_FUCHSIA)

    add_text(s, "Captions, hashtag stacks & Story rotation in 03-content-calendar.md & 05-social-templates.md",
             left=Inches(0.7), top=Inches(7.3), width=Inches(12), height=Inches(0.2),
             size=8, font=BODY, color=FG_LOW, align=PP_ALIGN.CENTER)

    set_notes(s, "Land KPI #5 (14/14 days), KPI #6 (4 formats — exceeds 3+) and KPI #9 setup right here. "
                 "The Planner screenshot is the rubric evidence; FB+IG screenshots show the live identities. "
                 "Time: 45 seconds.")


def slide_8_meta_ads(prs):
    s = add_blank_slide(prs)
    add_header(s, "Pillar 2 · Meta Ads", "$280 / 14 days · 4 ad sets · 6 creatives · all built to Review.", accent=ACCENT_CYAN)
    add_footer(s, "Pillar 2 · Meta Ads")

    # 4 actual Ads Manager Review screenshots in 2x2 grid
    sets = [
        ("Ad Set 1 · Students", "meta-ad-set-1-students.png", ACCENT_VIOLET),
        ("Ad Set 2 · Creators", "meta-ad-set-2-creators.png", ACCENT_FUCHSIA),
        ("Ad Set 3 · Knowledge Workers", "meta-ad-set-3-knowledge.png", ACCENT_CYAN),
        ("Ad Set 4 · Retargeting", "meta-ad-set-4-retarget.png", ACCENT_EMERALD),
    ]

    cell_w = Inches(5.85)
    cell_h = Inches(2.15)
    label_h = Inches(0.3)

    for i, (label, fname, accent) in enumerate(sets):
        col = i % 2
        row = i // 2
        x = Inches(0.7 + col * 6.2)
        y = Inches(2.45 + row * (cell_h + label_h + Inches(0.15) / Inches(1)))
        # Compute y manually so it's consistent
        y = Inches(2.45 + row * 2.55)

        add_text(s, label,
                 left=x, top=y, width=cell_w, height=label_h,
                 size=11, bold=True, font=BODY, color=accent)
        add_image(s, asset(fname),
                  left=x, top=y + label_h,
                  width=cell_w, height=cell_h,
                  border_color=accent, border_pt=1)

    add_text(s, "Built in Ads Manager up to Review · saved as Draft · never published. $0 actual spend.",
             left=Inches(0.7), top=Inches(7.0), width=Inches(12), height=Inches(0.3),
             size=10, font=BODY, color=FG_LOW, align=PP_ALIGN.CENTER)

    set_notes(s, "Showcase mode: 4 ad sets each with full Review-page screenshot. KPI #7 (Meta Ads "
                 "Strategy Completeness) and KPI #8 (Audience Definition Quality) — both Strong. Time: 55 seconds.")


def slide_9_conversation(prs):
    s = add_blank_slide(prs)
    add_header(s, "Pillar 2 · Always-on", "Welcome · away · FAQ · saved replies — KPI #9 evidence.", accent=ACCENT_CYAN)
    add_footer(s, "Pillar 2 · Templates")

    # 4 real screenshots from Business Suite Inbox / Automations
    items = [
        ("Greeting · DM welcome",       "auto-greeting.png",   ACCENT_VIOLET),
        ("Away message · outside hours", "auto-away.png",       ACCENT_FUCHSIA),
        ("FAQ shortcuts (5 set up)",    "auto-faq.png",        ACCENT_CYAN),
        ("Saved replies (5 templates)", "saved-replies.png",   ACCENT_EMERALD),
    ]

    cell_w = Inches(5.85)
    cell_h = Inches(2.15)

    for i, (label, fname, accent) in enumerate(items):
        col = i % 2
        row = i // 2
        x = Inches(0.7 + col * 6.2)
        y = Inches(2.45 + row * 2.55)

        add_text(s, label,
                 left=x, top=y, width=cell_w, height=Inches(0.3),
                 size=11, bold=True, font=BODY, color=accent)
        add_image(s, asset(fname),
                  left=x, top=y + Inches(0.3),
                  width=cell_w, height=cell_h,
                  border_color=accent, border_pt=1)

    add_text(s, "All four screenshots above = full evidence trail for KPI #9 (Automated Message / Welcome Note).",
             left=Inches(0.7), top=Inches(7.0), width=Inches(12), height=Inches(0.3),
             size=10, font=BODY, color=FG_LOW, align=PP_ALIGN.CENTER)

    set_notes(s, "All four screens come from Business Suite → Inbox → Automated Responses + Saved Replies. "
                 "Show them quickly, emphasise the always-on conversational layer means leads never wait. Time: 30 seconds.")


def _kpi_color(score: int):
    if score >= 4:
        return ACCENT_EMERALD
    if score == 3:
        return ACCENT_AMBER
    return RGBColor(0xF4, 0x3F, 0x5E)


def slide_10_keywords(prs):
    s = add_blank_slide(prs)
    add_header(s, "Pillar 4 · SEO", "Full 18-keyword strategy · long-tail-first · per page mapping.", accent=ACCENT_AMBER)
    add_footer(s, "Pillar 4 · Keyword research")

    # ── Top-row stats (compact)
    stats = [
        ("18", "total KWs", ACCENT_AMBER),
        ("5", "short-tail",  ACCENT_FUCHSIA),
        ("13", "long-tail",   ACCENT_EMERALD),
        ("21.8", "avg KD",      ACCENT_CYAN),
    ]
    for i, (n, lbl, c) in enumerate(stats):
        x = Inches(0.5 + i * 1.55)
        add_card(s, left=x, top=Inches(2.32), width=Inches(1.4), height=Inches(0.85), accent=c)
        add_text(s, n, left=x, top=Inches(2.4), width=Inches(1.4), height=Inches(0.4),
                 size=22, bold=True, font=DISPLAY, color=c, align=PP_ALIGN.CENTER)
        add_text(s, lbl.upper(), left=x, top=Inches(2.85), width=Inches(1.4), height=Inches(0.3),
                 size=8, font=BODY, color=FG_MID, align=PP_ALIGN.CENTER)

    # ── Full keyword table (18 rows) ─────────────────────────────────
    kws = [
        # (#, keyword, type, MSV, KD, CPC, intent, page, in_ads)
        (1,  "ai video summarizer",                          "Short · Head", "27,100", 48, "1.85", "Commercial",        "/",        True),
        (2,  "youtube video summarizer",                     "Short · Head", "22,200", 42, "1.40", "Commercial",        "/",        True),
        (3,  "summarize youtube video",                      "Short",        "18,100", 39, "1.20", "Transactional",     "/analyze", True),
        (4,  "ai for video transcripts",                     "Short",         "4,400", 31, "0.90", "Informational",     "/analyze", True),
        (5,  "live stream ai summary",                       "Short · Niche", "1,600", 24, "1.10", "Commercial",        "/live",    True),
        (6,  "summarize a 2 hour youtube video",             "Long-tail",     "2,400", 19, "0.75", "Transactional",     "/analyze", True),
        (7,  "how to summarize a long youtube lecture free", "Long-tail",       "880", 12, "0.55", "Informational",     "/analyze", True),
        (8,  "best free ai tool to summarize video",         "Long-tail",     "1,900", 22, "1.05", "Commercial",        "/",        True),
        (9,  "chat with youtube video ai",                   "Long-tail",     "1,300", 18, "1.30", "Transactional",     "/analyze", True),
        (10, "extract key moments from video ai",            "Long-tail",       "720", 14, "0.85", "Commercial",        "/analyze", True),
        (11, "ai timestamp summary youtube",                 "Long-tail",       "590", 11, "0.65", "Informational",     "/analyze", False),
        (12, "youtube live stream transcript real time",     "Long-tail · Niche","480", 16, "0.95", "Commercial",        "/live",    True),
        (13, "webinar to text ai summary",                   "Long-tail · B2B","1,100", 25, "1.45", "Commercial",        "/live",    True),
        (14, "study with youtube ai summary",                "Long-tail",       "720", 13, "0.70", "Students",          "/",        True),
        (15, "summarize trading livestream ai",              "Long-tail · Niche","210",  8, "0.60", "Commercial",        "/live",    False),
        (16, "medical lecture video summarizer",             "Long-tail · Niche","320", 14, "1.20", "Commercial",        "/analyze", False),
        (17, "ai video keyframe extraction",                 "Long-tail · Tech","480", 17, "0.95", "Informational",     "/",        False),
        (18, "ai video to notes converter",                  "Long-tail",     "1,600", 20, "0.90", "Transactional",     "/analyze", True),
    ]

    table_y = Inches(3.3)
    cols_w = [Inches(0.4), Inches(3.5), Inches(1.6), Inches(0.95), Inches(0.75), Inches(0.7), Inches(1.45), Inches(1.2), Inches(0.45)]
    cols_x = [Inches(0.5)]
    for w in cols_w[:-1]:
        cols_x.append(cols_x[-1] + w)

    headers = ["#", "Keyword", "Type", "MSV", "KD", "CPC", "Intent", "Page", "Ads"]
    # header row
    for i, h in enumerate(headers):
        align = PP_ALIGN.LEFT if i in (1, 2, 6, 7) else PP_ALIGN.CENTER
        add_text(s, h.upper(), left=cols_x[i], top=table_y, width=cols_w[i], height=Inches(0.28),
                 size=8, bold=True, font=BODY, color=FG_LOW, align=align)
    # divider
    div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.5), table_y + Inches(0.32),
                             Inches(11.05), Inches(0.015))
    div.fill.solid(); div.fill.fore_color.rgb = ACCENT_AMBER; div.line.fill.background()

    # data rows
    row_h = Inches(0.21)
    for ri, (n, kw, typ, msv, kd, cpc, intent, page, in_ads) in enumerate(kws):
        y = table_y + Inches(0.42) + row_h * ri
        # alternating row tint
        if ri % 2 == 0:
            tint = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(0.5), y - Inches(0.02),
                                      Inches(11.05), row_h)
            tint.fill.solid(); tint.fill.fore_color.rgb = RGBColor(0x14, 0x0E, 0x1E)
            tint.line.fill.background()

        kd_color = ACCENT_EMERALD if kd < 15 else (ACCENT_AMBER if kd < 30 else RGBColor(0xF4, 0x3F, 0x5E))

        add_text(s, str(n), left=cols_x[0], top=y, width=cols_w[0], height=row_h,
                 size=8, font=MONO, color=FG_LOW, align=PP_ALIGN.CENTER)
        add_text(s, kw, left=cols_x[1], top=y, width=cols_w[1], height=row_h,
                 size=9, font=BODY, color=FG_HIGH)
        add_text(s, typ, left=cols_x[2], top=y, width=cols_w[2], height=row_h,
                 size=8, font=BODY, color=FG_MID)
        add_text(s, msv, left=cols_x[3], top=y, width=cols_w[3], height=row_h,
                 size=9, font=MONO, color=ACCENT_CYAN, align=PP_ALIGN.RIGHT)
        add_text(s, str(kd), left=cols_x[4], top=y, width=cols_w[4], height=row_h,
                 size=9, bold=True, font=MONO, color=kd_color, align=PP_ALIGN.RIGHT)
        add_text(s, "$" + cpc, left=cols_x[5], top=y, width=cols_w[5], height=row_h,
                 size=9, font=MONO, color=FG_MID, align=PP_ALIGN.RIGHT)
        add_text(s, intent, left=cols_x[6], top=y, width=cols_w[6], height=row_h,
                 size=8, font=BODY, color=FG_MID)
        add_text(s, page, left=cols_x[7], top=y, width=cols_w[7], height=row_h,
                 size=8, font=MONO, color=ACCENT_VIOLET)
        add_text(s, "✓" if in_ads else "—", left=cols_x[8], top=y, width=cols_w[8], height=row_h,
                 size=10, bold=True, font=BODY,
                 color=ACCENT_EMERALD if in_ads else FG_LOW, align=PP_ALIGN.CENTER)

    add_text(s, "KD<15 (green) · KD 15-29 (amber) · KD 30+ (red)  ·  ✓ = also targeted in Google Ads (`08-google-ads-plan.md`)",
             left=Inches(0.5), top=Inches(7.1), width=Inches(12.5), height=Inches(0.25),
             size=8, font=BODY, color=FG_LOW, align=PP_ALIGN.CENTER)

    set_notes(s, "All 18 keywords with MSV, KD, CPC, intent, mapped page, and Google Ads inclusion. "
                 "Highlight: 5 short-tail head terms + 13 long-tail (meets KPI #15). "
                 "Long-tail-first because DA=1 makes head terms unrankable in 14 days. Time: 60 seconds.")


def slide_11_seo(prs):
    s = add_blank_slide(prs)
    add_header(s, "Pillar 4 · On-page SEO",
               "8 techniques applied · meta tags → Core Web Vitals · live evidence in GSC.",
               accent=ACCENT_AMBER)
    add_footer(s, "Pillar 4 · On-page SEO")

    # 8 techniques in a 4-col × 2-row grid · each card shows technique + body + measurable benefit.
    techniques = [
        ("Meta tags",
         "Per-page title + description + canonical + OG + Twitter card",
         "→ +CTR in SERP, social previews",
         ACCENT_VIOLET),
        ("Alt text",
         "Every image has descriptive alt — auto-generated from caption",
         "→ Image-search reach + a11y",
         ACCENT_CYAN),
        ("Header hierarchy",
         "Strict h1→h2→h3 — verified via Wave + axe DevTools",
         "→ Crawler topic-modelling",
         ACCENT_FUCHSIA),
        ("Structured data",
         "JSON-LD: Organization + WebSite + SoftwareApplication",
         "→ Rich-result eligibility",
         ACCENT_EMERALD),
        ("Sitemap + robots",
         "/sitemap.xml (4 URLs) + robots.txt — submitted to GSC",
         "→ Faster discovery / indexing",
         ACCENT_AMBER),
        ("Internal linking",
         "Hub-and-spoke: / → /analyze, /live, /library + cross-links",
         "→ Distributes link equity",
         ACCENT_VIOLET),
        ("Core Web Vitals",
         "next/image, lazy-load, Vercel Edge cache, font preload",
         "→ LCP < 2.5 s · CLS < 0.1",
         ACCENT_CYAN),
        ("Mobile-first + HTTPS",
         "Tailwind responsive · viewport meta · TLS via Vercel",
         "→ Mobile-index ranking signal",
         ACCENT_FUCHSIA),
    ]

    cols = 4
    cw = Inches(2.95); ch = Inches(1.55); gx = Inches(0.16); gy = Inches(0.18)
    sx = Inches(0.7); sy = Inches(2.35)

    for i, (title, body, benefit, c) in enumerate(techniques):
        col = i % cols
        row = i // cols
        x = sx + (cw + gx) * col
        y = sy + (ch + gy) * row
        add_card(s, left=x, top=y, width=cw, height=ch, accent=c)
        add_text(s, "✓",
                 left=x + Inches(0.15), top=y + Inches(0.1), width=Inches(0.4), height=Inches(0.4),
                 size=18, bold=True, font=DISPLAY, color=c)
        add_text(s, title,
                 left=x + Inches(0.6), top=y + Inches(0.1), width=cw - Inches(0.7), height=Inches(0.32),
                 size=12, bold=True, font=DISPLAY, color=FG_HIGH)
        add_text(s, body,
                 left=x + Inches(0.18), top=y + Inches(0.5), width=cw - Inches(0.36), height=Inches(0.65),
                 size=8, font=BODY, color=FG_MID, line_spacing=1.3)
        add_text(s, benefit,
                 left=x + Inches(0.18), top=y + Inches(1.18), width=cw - Inches(0.36), height=Inches(0.3),
                 size=8, bold=True, font=MONO, color=c)

    # Proofs row — Lighthouse + live evidence
    add_card(s, left=Inches(0.7), top=Inches(5.8), width=Inches(11.95), height=Inches(0.7),
             accent=ACCENT_VIOLET)
    add_text(s, "PROOFS",
             left=Inches(0.95), top=Inches(5.92), width=Inches(1.0), height=Inches(0.3),
             size=10, bold=True, font=BODY, color=ACCENT_VIOLET)
    add_text(s, "Lighthouse · SEO 100 / Best-Practices ≥ 90 / Performance audited mobile + desktop   ·   GSC: property verified · sitemap submitted · 3 URLs in priority crawl queue",
             left=Inches(2.0), top=Inches(5.92), width=Inches(10.4), height=Inches(0.5),
             size=9, font=BODY, color=FG_HIGH, line_spacing=1.35)

    # Live URL row
    add_text(s, "Live:  vidiq-two.vercel.app/sitemap.xml   ·   Robots:  /robots.txt   ·   Schema:  view-source: home   ·   Search Console: verified property",
             left=Inches(0.7), top=Inches(6.65), width=Inches(12), height=Inches(0.3),
             size=9, font=MONO, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

    # Aggregate benefit footer — "what does this earn us"
    add_text(s, "Net benefit: rank-eligible from day 1 · indexable to Google + Bing · social previews load with brand · WCAG 2.1 AA accessible · ~95th-percentile mobile UX.",
             left=Inches(0.7), top=Inches(7.0), width=Inches(12), height=Inches(0.3),
             size=9, font=BODY, color=FG_LOW, align=PP_ALIGN.CENTER, line_spacing=1.3)

    set_notes(s, "8 on-page SEO techniques — original 4 (meta, alt, headers, schema) plus 4 more "
                 "(sitemap+robots, internal linking, Core Web Vitals, mobile-first+HTTPS). Each card "
                 "names the technique, what we did, and the measurable benefit. Proofs row points to "
                 "live URLs the examiner can open: /sitemap.xml, /robots.txt, view-source for JSON-LD, "
                 "and GSC verification. Lighthouse run on demo morning. Time: 60 seconds.")


def slide_12_google_ads(prs):
    s = add_blank_slide(prs)
    add_header(s, "Pillar 4 · Google Ads", "$210 / 14 days · 3 campaigns · all built to Review.", accent=ACCENT_AMBER)
    add_footer(s, "Pillar 4 · Google Ads")

    campaigns = [
        ("Search · Brand Cluster",      "Maximise Conversions · 4 ad groups", "$70 / 14d",  "google-search-review.jpeg",  ACCENT_VIOLET),
        ("YouTube · In-Stream",         "Target CPM $4 · 30s skippable + 6s",  "$56 / 14d",  "google-youtube-review.jpeg", ACCENT_CYAN),
        ("Performance Max",              "Maximise Conversion Value · 3 asset groups", "$84 / 14d", "google-pmax-review.jpeg", ACCENT_FUCHSIA),
    ]

    cell_w = Inches(4.0)
    cell_h = Inches(2.6)

    for i, (name, sub, budget, fname, c) in enumerate(campaigns):
        x = Inches(0.5 + i * 4.2)

        # Title strip
        add_text(s, name,
                 left=x, top=Inches(2.4), width=cell_w, height=Inches(0.35),
                 size=14, bold=True, font=DISPLAY, color=c)
        add_text(s, sub,
                 left=x, top=Inches(2.75), width=cell_w, height=Inches(0.3),
                 size=10, font=BODY, color=FG_MID)
        add_text(s, budget,
                 left=x, top=Inches(3.05), width=cell_w, height=Inches(0.25),
                 size=9, font=MONO, color=c)

        # Real screenshot
        add_image(s, asset(fname, fname.replace(".jpeg", ".png")),
                  left=x, top=Inches(3.4),
                  width=cell_w, height=cell_h,
                  border_color=c, border_pt=1)

    # Bidding rules below
    add_text(s, "BIDDING + GUARDRAILS",
             left=Inches(0.5), top=Inches(6.2), width=Inches(5), height=Inches(0.25),
             size=9, bold=True, font=BODY, color=FG_LOW)
    rules = [
        "Start: Maximise Conversions  →  Switch to Target CPA $1.50 after 30 conv banked.",
        "Friday negative-keyword sweep: drop terms with > $5 spend, 0 conversions.",
        "PMax gated by conv volume; spare reallocates to Search if it doesn't unlock.",
    ]
    for i, r in enumerate(rules):
        add_text(s, f"·  {r}",
                 left=Inches(0.5), top=Inches(6.5 + i * 0.27), width=Inches(12.5), height=Inches(0.25),
                 size=9, font=BODY, color=FG_MID)

    set_notes(s, "All three screenshots are from the Review/Summary page in Google Ads (saved as Draft, "
                 "never published). Strategy from 08-google-ads-plan.md. KPI #17 (Google Ads Strategy "
                 "Completeness) and KPI #18 (Brand Alignment) — both Strong. Time: 50 seconds.")


def slide_13_budget(prs):
    s = add_blank_slide(prs)
    add_header(s, "Pillar 5 · Budget",
               "₨ 250 000 PKR · ~$890 USD · 14-day launch flight · ₨ 0 actual (showcase).",
               accent=ACCENT_EMERALD)
    add_footer(s, "Pillar 5 · Budget · PKR-primary (Pakistan launch)")

    # Six categories — PKR primary + USD secondary, same 38/28/16/8/7/3 weighting.
    lines = [
        ("Meta Ads",            "₨ 95,000",  "$340", "38%", ACCENT_VIOLET),
        ("Google Ads",          "₨ 70,000",  "$250", "28%", ACCENT_FUCHSIA),
        ("Creative production", "₨ 40,000",  "$143", "16%", ACCENT_CYAN),
        ("Influencer seeding",  "₨ 20,000",   "$71",  "8%", ACCENT_EMERALD),
        ("Tools & SaaS",        "₨ 17,500",   "$63",  "7%", ACCENT_AMBER),
        ("Contingency",         "₨ 7,500",    "$27",  "3%", RGBColor(0xF4, 0x3F, 0x5E)),
    ]
    for i, (name, pkr, usd, pct, c) in enumerate(lines):
        y = Inches(2.4 + i * 0.55)
        # Bar
        bar_w_total = Inches(5.6)
        pct_val = float(pct.replace("%", "")) / 100
        bar_w = Inches(5.6 * pct_val)
        bg_bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(2.85), y + Inches(0.1), bar_w_total, Inches(0.32))
        bg_bar.fill.solid()
        bg_bar.fill.fore_color.rgb = RGBColor(0x18, 0x10, 0x24)
        bg_bar.line.fill.background()
        bg_bar.adjustments[0] = 0.5
        fg_bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(2.85), y + Inches(0.1), bar_w, Inches(0.32))
        fg_bar.fill.solid()
        fg_bar.fill.fore_color.rgb = c
        fg_bar.line.fill.background()
        fg_bar.adjustments[0] = 0.5
        add_text(s, name,
                 left=Inches(0.7), top=y + Inches(0.1), width=Inches(2.1), height=Inches(0.4),
                 size=12, bold=True, font=DISPLAY, color=FG_HIGH)
        add_text(s, pkr,
                 left=Inches(8.55), top=y + Inches(0.1), width=Inches(1.55), height=Inches(0.4),
                 size=12, font=MONO, color=c, bold=True, align=PP_ALIGN.RIGHT)
        add_text(s, usd,
                 left=Inches(10.15), top=y + Inches(0.12), width=Inches(0.9), height=Inches(0.4),
                 size=10, font=MONO, color=FG_LOW, align=PP_ALIGN.RIGHT)
        add_text(s, pct,
                 left=Inches(11.1), top=y + Inches(0.1), width=Inches(0.9), height=Inches(0.4),
                 size=12, font=MONO, color=FG_MID, align=PP_ALIGN.RIGHT)

    # Why PKR 250,000? — rationale card
    add_card(s, left=Inches(0.7), top=Inches(5.75), width=Inches(11.95), height=Inches(1.15),
             accent=ACCENT_AMBER)
    add_text(s, "WHY ₨ 250,000?",
             left=Inches(0.95), top=Inches(5.85), width=Inches(3), height=Inches(0.3),
             size=10, bold=True, font=BODY, color=ACCENT_AMBER)
    rationale = (
        "Pakistani SaaS launch budgets cluster at ₨ 150 K – ₨ 400 K / month — ₨ 250 K is mid-range, aggressive enough "
        "for a 14-day flight.   ·   Meta CPM in Pakistan ≈ ₨ 80–250 (vs $5–15 in West) → 5–7× more reach per rupee.   ·   "
        "Google Search CPC for our long-tail KWs in PK ≈ ₨ 20–80 (vs $1–3 USD).   ·   ₨ 250 K buys ~80–120 K impressions "
        "across Meta + Google at our PK targeting — enough for the pixel to leave learning and converge to a CPA."
    )
    add_text(s, rationale,
             left=Inches(0.95), top=Inches(6.15), width=Inches(11.5), height=Inches(0.75),
             size=9, font=BODY, color=FG_MID, line_spacing=1.35)

    # Showcase actual
    add_card(s, left=Inches(0.7), top=Inches(7.0), width=Inches(11.95), height=Inches(0.45),
             accent=ACCENT_EMERALD)
    add_text(s, "SHOWCASE ACTUAL",
             left=Inches(0.95), top=Inches(7.05), width=Inches(2.7), height=Inches(0.35),
             size=10, bold=True, font=BODY, color=ACCENT_EMERALD)
    add_text(s, "₨ 0 spent — campaigns built to Review screen, screenshotted, never published. FX: 1 USD ≈ 280 PKR.",
             left=Inches(3.6), top=Inches(7.05), width=Inches(8.9), height=Inches(0.35),
             size=10, font=BODY, color=FG_HIGH)

    set_notes(s, "Budget is now PKR-primary because we are operating in Pakistan and pay Meta/Google in PKR. "
                 "₨ 250 K = mid-range Pakistani SaaS launch budget. Why the 38/28 split: paid acquisition is dominant "
                 "because organic SEO is months out for a DA-1 domain. Creative + influencer is intentionally lean. "
                 "Pakistani CPM/CPC rates are 5-7× cheaper than Western markets, so this PKR budget converts to "
                 "Western-equivalent reach of ~$3-4K spend. 09-budget.md has full rationale. Time: 60 seconds.")


def _3col_table(s, *, top, headers, rows, col_widths=None,
                left=Inches(0.5), accent_per_col=None,
                row_label_w=Inches(2.4), row_h=Inches(0.32),
                first_col_label=True):
    """Render a 3-column comparison table with optional first label column.

    `headers` = ['', 'VidIQ', 'NoteGPT', 'Eightify']  (first is label)
    `rows`    = list of [label, vidiq_val, notegpt_val, eightify_val]
    `accent_per_col` = [violet, cyan, amber] colours for each entity column header
    """
    if accent_per_col is None:
        accent_per_col = [ACCENT_VIOLET, ACCENT_CYAN, ACCENT_AMBER]

    if col_widths is None:
        # row label + 3 entity cols
        ent_w = Inches(3.18)
        col_widths = [row_label_w, ent_w, ent_w, ent_w]

    # column header strip
    x = left
    for ci, h in enumerate(headers):
        if ci == 0:
            add_text(s, h.upper() if h else "", left=x, top=top,
                     width=col_widths[0], height=Inches(0.32),
                     size=8, bold=True, font=BODY, color=FG_LOW)
        else:
            cap = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     x + Inches(0.05), top,
                                     col_widths[ci] - Inches(0.1), Inches(0.32))
            cap.fill.solid()
            cap.fill.fore_color.rgb = accent_per_col[ci - 1]
            cap.line.fill.background()
            cap.adjustments[0] = 0.3
            tf = cap.text_frame
            tf.margin_top = Inches(0.04)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = h
            r.font.name = DISPLAY
            r.font.size = Pt(11)
            r.font.bold = True
            r.font.color.rgb = BG_DEEP
        x += col_widths[ci]

    # data rows
    cur_y = top + Inches(0.42)
    for ri, row in enumerate(rows):
        # row tint (alternating)
        if ri % 2 == 0:
            tint = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      left, cur_y - Inches(0.02),
                                      sum((w.inches for w in col_widths)) * 914400,
                                      row_h)
            tint.fill.solid(); tint.fill.fore_color.rgb = RGBColor(0x14, 0x0E, 0x1E)
            tint.line.fill.background()
        x = left
        for ci, val in enumerate(row):
            color = ACCENT_VIOLET if ci == 0 and first_col_label else FG_HIGH
            font = BODY
            size = 9
            if ci == 0:
                color = FG_LOW; font = BODY; size = 9
            add_text(s, str(val),
                     left=x + Inches(0.08), top=cur_y,
                     width=col_widths[ci] - Inches(0.16), height=row_h,
                     size=size, font=font, color=color)
            x += col_widths[ci]
        cur_y += row_h

    return cur_y


def slide_14_competitive(prs):
    """Section A — Company overview + Section C — Website/SEO."""
    s = add_blank_slide(prs)
    add_header(s, "Pillar 5 · Competitive", "Company overview · pricing · website & SEO.", accent=ACCENT_EMERALD)
    add_footer(s, "Pillar 5 · Competitive — Company & Web")

    # Section A — Company overview
    add_text(s, "SECTION A · COMPANY OVERVIEW",
             left=Inches(0.5), top=Inches(2.3), width=Inches(8), height=Inches(0.3),
             size=10, bold=True, font=BODY, color=ACCENT_EMERALD)

    company_rows = [
        ("Brand name",       "VidIQ",                                              "NoteGPT",                       "Eightify"),
        ("Industry",         "AI · Video summarisation SaaS",                      "AI · Video summarisation",      "Chrome extension first"),
        ("Target audience",  "Students · educators · creators · domain pros",      "Students · knowledge workers",  "Power-YouTube viewers"),
        ("Value prop",       "Multimodal + grounded citations",                    "Free YT + PDF summariser",      "YouTube summary in 8 sec"),
        ("Pricing model",    "Free · paid tier roadmap",                           "Free 5/day → $4.99/mo",         "Free 4/day → $4.99/mo"),
        ("Website",          "vidiq-two.vercel.app",                               "notegpt.io",                    "eightify.app"),
    ]
    end_y = _3col_table(s, top=Inches(2.65),
                        headers=["ATTRIBUTE", "VidIQ", "NoteGPT", "Eightify"],
                        rows=company_rows,
                        accent_per_col=[ACCENT_VIOLET, ACCENT_CYAN, ACCENT_AMBER],
                        row_label_w=Inches(1.8))

    # Section C — Website & SEO
    add_text(s, "SECTION C · WEBSITE & SEO BASICS",
             left=Inches(0.5), top=end_y + Inches(0.15), width=Inches(8), height=Inches(0.3),
             size=10, bold=True, font=BODY, color=ACCENT_EMERALD)

    web_rows = [
        ("Domain Authority",   "1 (new)",                  "~28",                          "~31"),
        ("Monthly traffic",    "0 (new)",                  "~85,000",                      "~140,000"),
        ("Top KW #1",          "summarize 2-hr YT video",  "youtube summarizer",           "youtube summary chrome ext"),
        ("Top KW #2",          "live stream ai summary",   "summarize youtube video",      "eightify (brand)"),
        ("Top KW #3",          "chat with youtube ai",     "youtube to text",              "youtube ai summary"),
    ]
    _3col_table(s, top=end_y + Inches(0.5),
                headers=["ATTRIBUTE", "VidIQ", "NoteGPT", "Eightify"],
                rows=web_rows,
                accent_per_col=[ACCENT_VIOLET, ACCENT_CYAN, ACCENT_AMBER],
                row_label_w=Inches(1.8))

    set_notes(s, "Section A frames who's playing in this space. Section C exposes the authority gap "
                 "(DA 1 vs ~30) — that's why long-tail KW strategy in slide 10 is correct. "
                 "Time: 50 seconds.")


def slide_15_competitive_social(prs):
    """Section B — Social media presence (FB + IG)."""
    s = add_blank_slide(prs)
    add_header(s, "Pillar 5 · Competitive", "Social presence · Facebook & Instagram side-by-side.", accent=ACCENT_EMERALD)
    add_footer(s, "Pillar 5 · Competitive — Social")

    # Facebook block
    add_text(s, "📘  FACEBOOK",
             left=Inches(0.5), top=Inches(2.3), width=Inches(6), height=Inches(0.3),
             size=11, bold=True, font=DISPLAY, color=ACCENT_VIOLET)

    fb_rows = [
        ("Page URL",           "fb.com/…id=61588939576479", "facebook.com/notegpt",  "facebook.com/eightifyapp"),
        ("Total followers",    "Launching",         "~3,100",                "~620"),
        ("Avg likes/post",     "Launching",         "~12",                   "~4"),
        ("Posting freq",       "7/wk planned",      "~3/wk",                 "<1/wk"),
        ("Meta Ad activity",   "Setup phase",       "✓ Active (Q4 2025)",     "✗ Inactive"),
        ("Tone",               "Confident, calm",   "Functional, feature-led","Casual, meme-adjacent"),
    ]
    end_y = _3col_table(s, top=Inches(2.65),
                       headers=["ATTRIBUTE", "VidIQ", "NoteGPT", "Eightify"],
                       rows=fb_rows,
                       accent_per_col=[ACCENT_VIOLET, ACCENT_CYAN, ACCENT_AMBER],
                       row_label_w=Inches(1.8),
                       row_h=Inches(0.3))

    # Instagram block
    add_text(s, "📸  INSTAGRAM",
             left=Inches(0.5), top=end_y + Inches(0.15), width=Inches(6), height=Inches(0.3),
             size=11, bold=True, font=DISPLAY, color=ACCENT_FUCHSIA)

    ig_rows = [
        ("Handle",             "@vidiq_official",   "@notegpt.official",      "@eightifyapp"),
        ("Profile URL",        "instagram.com/vidiq_official/", "instagram.com/notegpt.official/", "instagram.com/eightifyapp/"),
        ("Total followers",    "Launching",         "~5,800",                 "~1,900"),
        ("Avg likes/post",     "Launching",         "~80",                    "~25"),
        ("Posting freq",       "7/wk planned",      "~4/wk",                  "~1/wk"),
        ("Hashtags/post",      "8-12",              "~6",                     "~3"),
    ]
    _3col_table(s, top=end_y + Inches(0.5),
                headers=["ATTRIBUTE", "VidIQ", "NoteGPT", "Eightify"],
                rows=ig_rows,
                accent_per_col=[ACCENT_VIOLET, ACCENT_CYAN, ACCENT_AMBER],
                row_label_w=Inches(1.8),
                row_h=Inches(0.3))

    set_notes(s, "We're behind on raw follower numbers (zero — we just launched), but our planned cadence "
                 "(7/wk) is more aggressive than either competitor. NoteGPT actively spends on Meta ads — "
                 "we'll match that with our 4-ad-set plan in slide 8. Time: 50 seconds.")


def slide_16_swot(prs):
    """SWOT analysis — Strengths · Weaknesses · Opportunities · Threats."""
    s = add_blank_slide(prs)
    add_header(s, "Pillar 5 · SWOT", "Strengths · Weaknesses · Opportunities · Threats.", accent=ACCENT_EMERALD)
    add_footer(s, "Pillar 5 · SWOT analysis")

    # 2×2 SWOT grid
    quad_w = Inches(6.0)
    quad_h = Inches(2.5)

    quadrants = [
        # (icon, label, color, x, y, items)
        ("💪", "STRENGTHS", ACCENT_EMERALD, Inches(0.5), Inches(2.3), [
            "Multimodal AI: vision + audio + text fused (neither competitor)",
            "Live-stream pipeline — rare in this category",
            "Provider-agnostic backend with auto-failover (Gemini ↔ Groq ↔ OpenAI)",
            "Domain-aware modes (medical · legal · trading · education)",
            "Open-source-able product — potential dev-community pull",
        ]),
        ("⚠️", "WEAKNESSES", RGBColor(0xF4, 0x3F, 0x5E), Inches(6.85), Inches(2.3), [
            "Zero brand awareness · DA = 1",
            "Brand-name collision with vidiq.com (YouTube SEO tool)",
            "No mobile app yet (web-only)",
            "Small shipping team",
            "Demo blocked by YouTube's cloud-IP rate-limiting (HF Space)",
        ]),
        ("🎯", "OPPORTUNITIES", ACCENT_CYAN, Inches(0.5), Inches(5.0), [
            "Live-stream summarisation — neither competitor ships",
            "Vertical domains (medical / legal / trading) — first-mover",
            "Pakistani + South Asian education — low CPC, no localised competition",
            "Open-source release — fork the dev-tool audience",
            "Creator-first positioning — neither competitor markets to creators",
        ]),
        ("🚨", "THREATS", ACCENT_AMBER, Inches(6.85), Inches(5.0), [
            "YouTube native AI summaries (Google rolling out 2025)",
            "OpenAI / Anthropic native ChatGPT-with-video features",
            "Eightify's 280k Chrome extension users (switching cost)",
            "Free-tier API tightening (Gemini / Groq)",
            "Brand confusion with vidiq.com (YouTube SEO, ~10M visits/mo)",
        ]),
    ]

    for icon, label, color, x, y, items in quadrants:
        # quadrant card
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, quad_w, quad_h)
        card.fill.solid(); card.fill.fore_color.rgb = BG_CARD
        card.line.color.rgb = color; card.line.width = Pt(0.75)
        card.adjustments[0] = 0.04

        # header strip
        add_text(s, f"{icon}  {label}",
                 left=x + Inches(0.2), top=y + Inches(0.15),
                 width=quad_w - Inches(0.4), height=Inches(0.35),
                 size=12, bold=True, font=DISPLAY, color=color)

        # divider
        sep = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 x + Inches(0.2), y + Inches(0.55),
                                 quad_w - Inches(0.4), Inches(0.015))
        sep.fill.solid(); sep.fill.fore_color.rgb = color; sep.line.fill.background()

        # bullets
        for i, it in enumerate(items):
            add_text(s, f"·  {it}",
                     left=x + Inches(0.25), top=y + Inches(0.65 + i * 0.34),
                     width=quad_w - Inches(0.5), height=Inches(0.32),
                     size=10, font=BODY, color=FG_HIGH, line_spacing=1.2)

    set_notes(s, "Lead with strengths — live + multimodal + verticals are the durable moats. "
                 "Address the weakness (zero awareness) by tying it to the launch ad strategy. "
                 "Opportunities are where 'why now' lives. Threats are honest — YouTube's native summaries "
                 "are the biggest risk. Time: 60 seconds.")


def slide_17_kpi_tracker_full(prs):
    """Full 18-row KPI table grouped by pillar."""
    s = add_blank_slide(prs)
    add_header(s, "Pillar 5 · KPI Tracker", "Self-assessment across 18 KPIs · 5 pillars · 4.83 / 5 average.", accent=ACCENT_EMERALD)
    add_footer(s, "Pillar 5 · KPI tracker")

    # Summary stats top-right
    summary = [
        ("17", "STRONG", ACCENT_EMERALD),
        ("1", "ADEQUATE", ACCENT_AMBER),
        ("0", "NEEDS WORK", RGBColor(0xF4, 0x3F, 0x5E)),
        ("4.83", "AVG / 5", ACCENT_VIOLET),
    ]
    for i, (n, lbl, c) in enumerate(summary):
        x = Inches(0.5 + i * 1.55)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, Inches(2.3), Inches(1.4), Inches(0.85))
        card.fill.solid(); card.fill.fore_color.rgb = BG_CARD
        card.line.color.rgb = c; card.line.width = Pt(0.75)
        card.adjustments[0] = 0.1
        add_text(s, n, left=x, top=Inches(2.38), width=Inches(1.4), height=Inches(0.45),
                 size=22, bold=True, font=DISPLAY, color=c, align=PP_ALIGN.CENTER)
        add_text(s, lbl, left=x, top=Inches(2.85), width=Inches(1.4), height=Inches(0.3),
                 size=8, font=BODY, color=FG_MID, align=PP_ALIGN.CENTER)

    # 18 KPI rows
    kpis = [
        # (#, pillar, name, score, evidence_short)
        (1,  "Brand", "Logo Design Completeness",          5, "Vector SVG + PNG · typography + colour rationale"),
        (2,  "Brand", "Colour Psychology Justification",   5, "01-brand-guide §4 — HSL/Hex + Mehta-Zhu citation"),
        (3,  "Brand", "30-Sec Video Ad Quality",           3, "Script + storyboard + master cut filmed (slide 6)"),
        (4,  "Brand", "Brand Consistency",                 5, "Audit across 7 surfaces — favicon · nav · hero · OG"),
        (5,  "Social","Content Calendar Coverage",         5, "28 posts (14 days × 2 channels) — Planner ss"),
        (6,  "Social","Post Type Variety",                 5, "4 formats — Reel · Carousel · Single · Story"),
        (7,  "Social","Meta Ads Strategy Completeness",    5, "Objective ✓ budget ✓ 4 ad sets ✓ 6 creatives ✓"),
        (8,  "Social","Target Audience Definition",        5, "Each ad set: age · location · interests · behaviour"),
        (9,  "Social","Welcome / Auto-Reply",              4, "Greeting · away · 5 FAQs · 5 saved replies"),
        (10, "Product","Functionality (0 broken links)",   5, "Live deploy · 9 routes 200 OK · vidiq-two.vercel.app"),
        (11, "Product","Brand Identity on Website",        5, "7 surfaces audited · same palette throughout"),
        (12, "Product","UI/UX Clarity",                    5, "Single CTA · TopNav · Aurora hero · TanStack Query"),
        (13, "Product","Mobile Responsiveness",            5, "Tailwind responsive · viewport · 390 px verified"),
        (14, "SEO",   "Keyword Research Depth",            5, "18 keywords with MSV + KD + CPC (slide 10)"),
        (15, "SEO",   "Long-tail vs Short-tail Balance",   5, "5 short-tail + 13 long-tail (≥5 each)"),
        (16, "SEO",   "On-Page SEO Coverage",              5, "Meta tags ✓ alt text ✓ headers ✓ KW usage ✓"),
        (17, "SEO",   "Google Ads Strategy Completeness",  5, "3 campaigns · 4 ad groups · RSAs · sitelinks (slide 12)"),
        (18, "SEO",   "Ad Creative Brand Alignment",       5, "All ad creative uses Pillar 1 visuals (slides 5-8)"),
    ]

    PILLAR_COLOR = {
        "Brand":   ACCENT_VIOLET,
        "Social":  ACCENT_CYAN,
        "Product": ACCENT_EMERALD,
        "SEO":     ACCENT_AMBER,
    }

    table_y = Inches(3.4)
    cols_w = [Inches(0.4), Inches(0.85), Inches(3.6), Inches(0.65), Inches(5.95)]
    cols_x = [Inches(0.5)]
    for w in cols_w[:-1]:
        cols_x.append(cols_x[-1] + w)

    # header
    headers = ["#", "PILLAR", "KPI METRIC", "SCORE", "EVIDENCE"]
    for i, h in enumerate(headers):
        add_text(s, h, left=cols_x[i], top=table_y, width=cols_w[i], height=Inches(0.26),
                 size=8, bold=True, font=BODY, color=FG_LOW)
    # divider
    div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.5), table_y + Inches(0.3),
                             Inches(11.45), Inches(0.015))
    div.fill.solid(); div.fill.fore_color.rgb = ACCENT_EMERALD
    div.line.fill.background()

    row_h = Inches(0.18)
    for ri, (n, pillar, name, score, evidence) in enumerate(kpis):
        y = table_y + Inches(0.4) + row_h * ri
        if ri % 2 == 0:
            tint = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(0.5), y - Inches(0.02),
                                      Inches(11.45), row_h)
            tint.fill.solid(); tint.fill.fore_color.rgb = RGBColor(0x14, 0x0E, 0x1E)
            tint.line.fill.background()

        score_color = _kpi_color(score)
        pillar_color = PILLAR_COLOR.get(pillar, ACCENT_VIOLET)

        add_text(s, str(n), left=cols_x[0], top=y, width=cols_w[0], height=row_h,
                 size=8, font=MONO, color=FG_LOW, align=PP_ALIGN.CENTER)
        # pillar pill
        pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  cols_x[1] + Inches(0.04), y + Inches(0.005),
                                  cols_w[1] - Inches(0.12), Inches(0.16))
        pill.fill.solid(); pill.fill.fore_color.rgb = pillar_color
        pill.line.fill.background()
        pill.adjustments[0] = 0.5
        tf = pill.text_frame
        tf.margin_left = Inches(0.04); tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = pillar.upper(); r.font.name = BODY; r.font.size = Pt(7); r.font.bold = True
        r.font.color.rgb = BG_DEEP

        add_text(s, name, left=cols_x[2], top=y, width=cols_w[2], height=row_h,
                 size=8, font=BODY, color=FG_HIGH)
        add_text(s, f"{score} / 5", left=cols_x[3], top=y, width=cols_w[3], height=row_h,
                 size=9, bold=True, font=DISPLAY, color=score_color, align=PP_ALIGN.CENTER)
        add_text(s, evidence, left=cols_x[4], top=y, width=cols_w[4], height=row_h,
                 size=7, font=BODY, color=FG_MID)

    add_text(s, "KPI #3 (video ad) at 3 — closes to 5 once edited cutdowns are exported.",
             left=Inches(0.5), top=Inches(7.05), width=Inches(12), height=Inches(0.25),
             size=9, font=BODY, color=FG_LOW)

    set_notes(s, "All 18 KPIs scored, evidence cited, average 4.83/5. The xlsx in data/DM_Competitive_KPI.xlsx "
                 "has the same scores. The single Adequate (KPI #3) closes to Strong once the video ad is "
                 "fully edited. Time: 60 seconds.")


def slide_18_kpi_scorecard(prs):
    s = add_blank_slide(prs)
    add_header(s, "Pillar 5 · KPI scorecard", "17 / 18 strong  ·  4.83 / 5 average.", accent=ACCENT_EMERALD)
    add_footer(s, "Pillar 5 · KPI scorecard")

    # 6×3 grid (18 KPIs) with status tint
    KPIS = [
        # (number, score, status_color)
        (1, 5, ACCENT_EMERALD), (2, 5, ACCENT_EMERALD), (3, 3, ACCENT_AMBER),
        (4, 5, ACCENT_EMERALD), (5, 5, ACCENT_EMERALD), (6, 5, ACCENT_EMERALD),
        (7, 5, ACCENT_EMERALD), (8, 5, ACCENT_EMERALD), (9, 4, ACCENT_EMERALD),
        (10, 5, ACCENT_EMERALD), (11, 5, ACCENT_EMERALD), (12, 5, ACCENT_EMERALD),
        (13, 5, ACCENT_EMERALD), (14, 5, ACCENT_EMERALD), (15, 5, ACCENT_EMERALD),
        (16, 5, ACCENT_EMERALD), (17, 5, ACCENT_EMERALD), (18, 5, ACCENT_EMERALD),
    ]
    cell_w = Inches(2.0)
    cell_h = Inches(0.7)
    gap = Inches(0.1)
    cols = 6
    rows = 3
    total_w = cell_w * cols + gap * (cols - 1)
    start_x = (SLIDE_W - total_w) / 2
    start_y = Inches(2.6)
    for i, (n, score, c) in enumerate(KPIS):
        col = i % cols
        row = i // cols
        x = start_x + (cell_w + gap) * col
        y = start_y + (cell_h + gap) * row
        cell = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cell_w, cell_h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BG_CARD
        cell.line.color.rgb = c
        cell.line.width = Pt(1.5)
        cell.adjustments[0] = 0.18
        add_text(s, f"KPI {n:02d}",
                 left=x, top=y + Inches(0.1), width=cell_w, height=Inches(0.25),
                 size=10, font=MONO, color=FG_LOW, align=PP_ALIGN.CENTER)
        add_text(s, f"{score} / 5",
                 left=x, top=y + Inches(0.32), width=cell_w, height=Inches(0.4),
                 size=18, bold=True, font=DISPLAY, color=c, align=PP_ALIGN.CENTER)

    # Footer note
    add_text(s, "1 KPI at 3 (KPI #3 — 30s video ad, awaiting filming).  Closes to 5 once shot.",
             left=Inches(0.7), top=Inches(6.5), width=Inches(12), height=Inches(0.4),
             size=11, font=BODY, color=FG_MID, align=PP_ALIGN.CENTER)

    set_notes(s, "Land the 17/18 number first, then explain the open one (video ad). "
                 "Cite the live KPI tracker in the app: vidiq-two.vercel.app/marketing. Time: 50 seconds.")


def slide_19_qa_backup(prs):
    s = add_blank_slide(prs)
    add_header(s, "Q&A · Backup", "How citations work — multimodal grounding.", accent=ACCENT_VIOLET)
    add_footer(s, "Backup · Q&A")

    add_text(s, "Q. \"How do you compete with YouTube's own AI summary?\"",
             left=Inches(0.7), top=Inches(2.5), width=Inches(12), height=Inches(0.5),
             size=18, bold=True, font=DISPLAY, color=FG_HIGH)
    add_text(s,
             "Three durable moats:",
             left=Inches(0.7), top=Inches(3.1), width=Inches(12), height=Inches(0.4),
             size=14, font=BODY, color=FG_MID)

    moats = [
        ("Live-stream support", "YouTube's auto-summary is post-hoc only. We process live streams in 30s windows."),
        ("Multimodal grounding", "Every claim cites the exact transcript moment + linked frame. YouTube's is generic."),
        ("Vertical-domain modes", "Medical / legal / trading prompts adapt the analysis to domain semantics."),
    ]
    for i, (title, body) in enumerate(moats):
        y = Inches(3.7 + i * 1.0)
        add_card(s, left=Inches(0.7), top=y, width=Inches(11.95), height=Inches(0.85), accent=ACCENT_VIOLET)
        add_text(s, f"{i+1}. {title}",
                 left=Inches(1), top=y + Inches(0.13), width=Inches(4), height=Inches(0.4),
                 size=14, bold=True, font=DISPLAY, color=ACCENT_VIOLET)
        add_text(s, body,
                 left=Inches(5.15), top=y + Inches(0.18), width=Inches(7.5), height=Inches(0.6),
                 size=12, font=BODY, color=FG_HIGH, line_spacing=1.3)

    set_notes(s, "Pull this slide out only if asked. Don't volunteer it. Time: 40 seconds if used.")


# ═══════════════════════════════════════════════════════════════════════════
# DIGITAL-MARKETING APPENDIX  ·  Slides 20-30
# Campaign-level detail · what we did · what the aim was · pulled from
# marketing/01..16-*.md plus assets/firefly + assets/screenshots.
# ═══════════════════════════════════════════════════════════════════════════


def slide_20_strategy_overview(prs):
    s = add_blank_slide(prs)
    add_header(s, "Marketing strategy",
               "Five pillars · one funnel · zero spend (showcase mode).",
               accent=ACCENT_VIOLET)
    add_footer(s, "Strategy · pillars · positioning")

    # Top row — three "thesis" cards
    thesis = [
        ("Mission",
         "Watch less. Learn more.",
         "Turn any video — recorded or live — into structured intelligence in seconds.",
         ACCENT_VIOLET),
        ("Positioning",
         "Multimodal · grounded · live",
         "Vision + audio + LLM with timestamp citations · live-stream pipeline · vertical-domain modes (medical · legal · trading).",
         ACCENT_FUCHSIA),
        ("GTM thesis",
         "Long-tail first · twin paid stack",
         "Win KD<15 KWs in 90 days · Meta + Google ads to seed conversion data · convert via free first analysis.",
         ACCENT_CYAN),
    ]
    card_w = Inches(4.05); card_h = Inches(2.0); gap = Inches(0.12)
    start_x = Inches(0.7); start_y = Inches(2.35)
    for i, (eyebrow, title, body, accent) in enumerate(thesis):
        x = start_x + (card_w + gap) * i
        add_card(s, left=x, top=start_y, width=card_w, height=card_h, accent=accent)
        add_text(s, eyebrow.upper(),
                 left=x + Inches(0.25), top=start_y + Inches(0.18),
                 width=card_w - Inches(0.5), height=Inches(0.25),
                 size=9, bold=True, font=BODY, color=accent)
        add_text(s, title,
                 left=x + Inches(0.25), top=start_y + Inches(0.45),
                 width=card_w - Inches(0.5), height=Inches(0.55),
                 size=18, bold=True, font=DISPLAY, color=FG_HIGH, line_spacing=1.1)
        add_text(s, body,
                 left=x + Inches(0.25), top=start_y + Inches(1.05),
                 width=card_w - Inches(0.5), height=Inches(0.95),
                 size=10, font=BODY, color=FG_MID, line_spacing=1.3)

    # Bottom row — five pillars chip strip
    pillars = [
        ("P1 · Brand",   "Logos · palette · 30 s ad",            ACCENT_VIOLET),
        ("P2 · Social",  "14-day calendar · Meta Ads · Inbox",   ACCENT_FUCHSIA),
        ("P3 · Product", "Live web app · 8 routes · 200 OK",     ACCENT_CYAN),
        ("P4 · SEO/SEM", "18 KWs · on-page audit · Google Ads",  ACCENT_EMERALD),
        ("P5 · KPI",     "$740 plan · 17/18 strong · 4.83 / 5",  ACCENT_AMBER),
    ]
    pw = Inches(2.42); ph = Inches(1.45); pg = Inches(0.05)
    sx = Inches(0.7); sy = Inches(4.6)
    for i, (title, body, accent) in enumerate(pillars):
        x = sx + (pw + pg) * i
        add_card(s, left=x, top=sy, width=pw, height=ph, accent=accent)
        add_text(s, title,
                 left=x + Inches(0.18), top=sy + Inches(0.18),
                 width=pw - Inches(0.36), height=Inches(0.4),
                 size=14, bold=True, font=DISPLAY, color=accent)
        add_text(s, body,
                 left=x + Inches(0.18), top=sy + Inches(0.6),
                 width=pw - Inches(0.36), height=Inches(0.8),
                 size=9, font=BODY, color=FG_HIGH, line_spacing=1.3)

    # Headline stats strip
    stats = [
        ("$0", "actual spend"),
        ("$740", "planned launch"),
        ("28 / 28", "posts queued"),
        ("4 + 3", "Meta + Google campaigns"),
        ("4.83 / 5", "KPI average"),
    ]
    sw = Inches(2.42); sh = Inches(0.55)
    for i, (big, label) in enumerate(stats):
        x = sx + (sw + pg) * i
        y = Inches(6.18)
        add_text(s, big,
                 left=x, top=y, width=sw, height=Inches(0.32),
                 size=18, bold=True, font=DISPLAY, color=ACCENT_VIOLET, align=PP_ALIGN.CENTER)
        add_text(s, label.upper(),
                 left=x, top=y + Inches(0.32), width=sw, height=Inches(0.22),
                 size=8, bold=True, font=BODY, color=FG_LOW, align=PP_ALIGN.CENTER)

    set_notes(s, "Open the appendix here. The five pillars map 1-to-1 onto the rubric. "
                 "Showcase-mode means we built every campaign to the *Review* step in Ads Manager "
                 "without publishing — that's the deliverable the brief asked for. Time: 60 s.")


def slide_21_brand_identity(prs):
    s = add_blank_slide(prs)
    add_header(s, "Pillar 1 · Brand identity & voice",
               "Why violet, why these fonts, why this voice.",
               accent=ACCENT_VIOLET)
    add_footer(s, "Pillar 1 · Brand · 01-brand-guide.md")

    # Left half — palette swatches + rationale
    add_text(s, "Palette",
             left=Inches(0.7), top=Inches(2.35), width=Inches(6), height=Inches(0.35),
             size=14, bold=True, font=DISPLAY, color=FG_HIGH)

    swatches = [
        (BG_DEEP,        "#0a0612", "App canvas",     "Cinema-feel near-black with violet undertone"),
        (ACCENT_VIOLET,  "#a855f7", "Brand violet",   "Trust (blue) + creativity (red) = insight"),
        (ACCENT_FUCHSIA, "#ec4899", "Gradient pair",  "Energy + AI cue (Linear, Midjourney)"),
        (ACCENT_CYAN,    "#06b6d4", "Live signal",    "Used sparingly for 'live' badges + stats"),
        (ACCENT_EMERALD, "#10b981", "Success",        "Calmer than Tailwind default green"),
    ]
    for i, (rgb, hex_, role, why) in enumerate(swatches):
        y = Inches(2.78 + i * 0.55)
        # color chip
        chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(0.7), y, Inches(0.55), Inches(0.45))
        chip.fill.solid(); chip.fill.fore_color.rgb = rgb
        chip.line.color.rgb = RGBColor(0x2A, 0x22, 0x3A); chip.line.width = Pt(0.5)
        chip.adjustments[0] = 0.3
        # hex
        add_text(s, hex_,
                 left=Inches(1.35), top=y + Inches(0.05), width=Inches(0.95), height=Inches(0.4),
                 size=10, font=MONO, color=FG_HIGH, bold=True)
        # role
        add_text(s, role,
                 left=Inches(2.35), top=y + Inches(0.04), width=Inches(1.6), height=Inches(0.4),
                 size=10, font=BODY, color=ACCENT_VIOLET, bold=True)
        # why
        add_text(s, why,
                 left=Inches(4.0), top=y + Inches(0.06), width=Inches(2.6), height=Inches(0.4),
                 size=9, font=BODY, color=FG_MID)

    # Right half — typography card
    add_text(s, "Typography",
             left=Inches(7.1), top=Inches(2.35), width=Inches(5.5), height=Inches(0.35),
             size=14, bold=True, font=DISPLAY, color=FG_HIGH)

    fonts = [
        ("Plus Jakarta Sans", "Display · headlines · hero",
         "Geometric, slightly humanist — modern at 40 pt without feeling austere.", DISPLAY, 22),
        ("Inter",              "Body · UI labels",
         "Most-tested SaaS body face — crisp at 14 px where dashboards live.", BODY, 16),
        ("JetBrains Mono",     "Timestamps · code",
         "Carries developer-tool credibility for the strategy → pseudocode feature.", MONO, 14),
    ]
    for i, (name, role, why, fontname, sample_pt) in enumerate(fonts):
        y = Inches(2.78 + i * 1.05)
        add_card(s, left=Inches(7.1), top=y, width=Inches(5.55), height=Inches(0.95),
                 accent=ACCENT_FUCHSIA)
        add_text(s, name,
                 left=Inches(7.3), top=y + Inches(0.1), width=Inches(3.5), height=Inches(0.35),
                 size=sample_pt, bold=True, font=fontname, color=FG_HIGH)
        add_text(s, role.upper(),
                 left=Inches(10.7), top=y + Inches(0.13), width=Inches(1.85), height=Inches(0.25),
                 size=8, bold=True, font=BODY, color=ACCENT_FUCHSIA, align=PP_ALIGN.RIGHT)
        add_text(s, why,
                 left=Inches(7.3), top=y + Inches(0.5), width=Inches(5.2), height=Inches(0.4),
                 size=9, font=BODY, color=FG_MID, line_spacing=1.3)

    # Bottom — voice rules (do/don't)
    add_text(s, "Voice rules",
             left=Inches(0.7), top=Inches(5.85), width=Inches(8), height=Inches(0.3),
             size=14, bold=True, font=DISPLAY, color=FG_HIGH)
    # Do card
    add_card(s, left=Inches(0.7), top=Inches(6.18), width=Inches(5.95), height=Inches(0.85),
             accent=ACCENT_EMERALD)
    add_text(s, "✓  DO",
             left=Inches(0.85), top=Inches(6.25), width=Inches(1), height=Inches(0.3),
             size=11, bold=True, font=BODY, color=ACCENT_EMERALD)
    add_text(s, "Lead with the verb · cite specifics · second person · short declarative sentences.",
             left=Inches(1.7), top=Inches(6.27), width=Inches(4.85), height=Inches(0.7),
             size=10, font=BODY, color=FG_HIGH, line_spacing=1.3)
    # Don't card
    add_card(s, left=Inches(6.75), top=Inches(6.18), width=Inches(5.9), height=Inches(0.85),
             accent=ACCENT_FUCHSIA)
    add_text(s, "✗  DON'T",
             left=Inches(6.9), top=Inches(6.25), width=Inches(1.2), height=Inches(0.3),
             size=11, bold=True, font=BODY, color=ACCENT_FUCHSIA)
    add_text(s, "\"revolutionary\" · \"100 % accurate\" · \"AI understands\" · long abstract claims.",
             left=Inches(7.95), top=Inches(6.27), width=Inches(4.55), height=Inches(0.7),
             size=10, font=BODY, color=FG_HIGH, line_spacing=1.3)

    set_notes(s, "Why violet not blue: every SaaS uses blue; violet adds the AI/creativity association "
                 "without breaking trust signal. Citation: Mehta-Zhu (2009), Labrecque-Milne (2012). "
                 "Three-font stack carries display, body, and monospace for timestamps. Time: 50 s.")


def slide_22_video_ad_campaign(prs):
    s = add_blank_slide(prs)
    add_header(s, "Pillar 1 · 30-second video ad",
               "Sizzle + value prop · 4 cutdowns from one shoot.",
               accent=ACCENT_VIOLET)
    add_footer(s, "Pillar 1 · Video ad · 02-video-ad-script.md")

    # Top brief card
    brief = [
        ("Goal",     "Drive sign-ups (/analyze) — emotional hook + clear value prop"),
        ("Audience", "Students cramming · creators researching · pros staring at 90-min webinars"),
        ("Promise",  "Watch the highlight reel of every video before you commit your hour"),
        ("Tone",     "Confident, calm, slightly cinematic — never frantic"),
    ]
    add_card(s, left=Inches(0.7), top=Inches(2.35), width=Inches(11.95), height=Inches(1.1),
             accent=ACCENT_VIOLET)
    for i, (k, v) in enumerate(brief):
        col = i % 2; row = i // 2
        x = Inches(0.95 + col * 5.95)
        y = Inches(2.45 + row * 0.45)
        add_text(s, k.upper(),
                 left=x, top=y, width=Inches(0.9), height=Inches(0.25),
                 size=9, bold=True, font=BODY, color=ACCENT_VIOLET)
        add_text(s, v,
                 left=x + Inches(0.95), top=y, width=Inches(4.85), height=Inches(0.35),
                 size=10, font=BODY, color=FG_HIGH, line_spacing=1.3)

    # Beat-by-beat strip — 6 cards in 2x3 grid
    beats = [
        ("0:00 – 0:03", "Phone close-up",      "YT bar 0:42 / 1:48:00 — \"you hit play on a 90-min video…\""),
        ("0:03 – 0:06", "Time-lapse dim",      "Day → night, hand still on lap — \"…and lose an evening.\""),
        ("0:06 – 0:10", "VidIQ Analyse form",  "URL paste · click — caption: \"What if you could read it instead?\""),
        ("0:10 – 0:18", "Beat montage (4×)",   "Transcript · Keyframes · Summary · Chat with citations"),
        ("0:18 – 0:23", "Workspace zoom-out",  "Click chapter → player seeks — \"every claim grounded in source\""),
        ("0:23 – 0:30", "Logo end-card",       "VidIQ · \"AI Video Intelligence\" · Watch less. Learn more."),
    ]
    bw = Inches(3.93); bh = Inches(1.4); bg = Inches(0.08)
    sx = Inches(0.7); sy = Inches(3.6)
    for i, (t, scene, caption) in enumerate(beats):
        col = i % 3; row = i // 3
        x = sx + (bw + bg) * col
        y = sy + (bh + bg) * row
        add_card(s, left=x, top=y, width=bw, height=bh, accent=ACCENT_FUCHSIA)
        add_text(s, t,
                 left=x + Inches(0.18), top=y + Inches(0.13),
                 width=Inches(2), height=Inches(0.25),
                 size=10, bold=True, font=MONO, color=ACCENT_FUCHSIA)
        add_text(s, scene,
                 left=x + Inches(0.18), top=y + Inches(0.42),
                 width=bw - Inches(0.36), height=Inches(0.35),
                 size=12, bold=True, font=DISPLAY, color=FG_HIGH)
        add_text(s, caption,
                 left=x + Inches(0.18), top=y + Inches(0.78),
                 width=bw - Inches(0.36), height=Inches(0.55),
                 size=9, font=BODY, color=FG_MID, line_spacing=1.3)

    # Cutdowns row
    cuts = [("30s", "9:16", "IG/TikTok/Shorts — master"),
            ("15s", "9:16", "Meta Ads — Reels"),
            ("6s",  "1:1",  "Feed bumper · YT non-skippable"),
            ("10s", "16:9", "YouTube pre-roll skippable")]
    cw = Inches(2.95); ch = Inches(0.55); cg = Inches(0.07)
    cx = Inches(0.7); cy = Inches(6.55)
    for i, (length, ar, use) in enumerate(cuts):
        x = cx + (cw + cg) * i
        add_card(s, left=x, top=cy, width=cw, height=ch, accent=ACCENT_CYAN)
        add_text(s, length,
                 left=x + Inches(0.15), top=cy + Inches(0.13),
                 width=Inches(0.8), height=Inches(0.3),
                 size=14, bold=True, font=DISPLAY, color=ACCENT_CYAN)
        add_text(s, ar,
                 left=x + Inches(0.85), top=cy + Inches(0.16),
                 width=Inches(0.6), height=Inches(0.25),
                 size=9, bold=True, font=MONO, color=FG_LOW)
        add_text(s, use,
                 left=x + Inches(1.4), top=cy + Inches(0.16),
                 width=cw - Inches(1.55), height=Inches(0.3),
                 size=9, font=BODY, color=FG_HIGH)

    set_notes(s, "Aim: drive sign-ups via emotional hook (lost evening) + crisp value prop. "
                 "Master 30 s plus three cutdowns from the same shoot — Meta Ads C1/C2/C3 plus a "
                 "10 s 16:9 for YouTube pre-roll. Music: cinematic-tech 95 BPM. End-card: wordmark + "
                 "tagline. Time: 60 s.")


def slide_23_content_calendar(prs):
    s = add_blank_slide(prs)
    add_header(s, "Pillar 2 · 14-day content calendar",
               "28 posts · 4 formats · Facebook + Instagram cross-posted.",
               accent=ACCENT_FUCHSIA)
    add_footer(s, "Pillar 2 · Calendar · 03-content-calendar.md")

    # Top — cadence chips
    chips = [
        ("Posts/wk",     "7 / channel",        ACCENT_FUCHSIA),
        ("Total flight", "28 posts (14d × 2)", ACCENT_VIOLET),
        ("Best times",   "10:30 + 20:30 PKT",  ACCENT_CYAN),
        ("Hashtags",     "8 – 12 / IG post",   ACCENT_EMERALD),
    ]
    cw = Inches(2.95); ch = Inches(0.7); cg = Inches(0.08)
    for i, (k, v, accent) in enumerate(chips):
        x = Inches(0.7) + (cw + cg) * i
        add_card(s, left=x, top=Inches(2.35), width=cw, height=ch, accent=accent)
        add_text(s, k.upper(),
                 left=x + Inches(0.18), top=Inches(2.43), width=Inches(1.6), height=Inches(0.25),
                 size=8, bold=True, font=BODY, color=accent)
        add_text(s, v,
                 left=x + Inches(0.18), top=Inches(2.65), width=cw - Inches(0.36), height=Inches(0.4),
                 size=14, bold=True, font=DISPLAY, color=FG_HIGH)

    # Content pillars — left card
    add_text(s, "Content pillars (rotation)",
             left=Inches(0.7), top=Inches(3.3), width=Inches(6), height=Inches(0.3),
             size=12, bold=True, font=DISPLAY, color=FG_HIGH)
    pillars = [
        ("Show product",    "P", 35, ACCENT_VIOLET),
        ("Use cases",       "U", 25, ACCENT_FUCHSIA),
        ("Educate",         "E", 20, ACCENT_CYAN),
        ("Social proof",    "S", 10, ACCENT_EMERALD),
        ("Lighten up",      "L", 10, ACCENT_AMBER),
    ]
    bar_x = Inches(0.7); bar_w = Inches(6.0); bar_y = Inches(3.7); bar_h = Inches(0.3)
    cum = 0.0
    for name, code, pct, accent in pillars:
        seg_w = Inches(bar_w.inches * pct / 100.0)
        seg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(bar_x.inches + bar_w.inches * cum / 100.0),
                                  bar_y, seg_w, bar_h)
        seg.fill.solid(); seg.fill.fore_color.rgb = accent
        seg.line.fill.background()
        cum += pct
    # Pillar legend rows
    for i, (name, code, pct, accent) in enumerate(pillars):
        y = Inches(4.15 + i * 0.32)
        # color chip
        chip = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.72), y + Inches(0.07),
                                   Inches(0.18), Inches(0.18))
        chip.fill.solid(); chip.fill.fore_color.rgb = accent
        chip.line.fill.background()
        add_text(s, name,
                 left=Inches(1.0), top=y + Inches(0.04), width=Inches(3.5), height=Inches(0.3),
                 size=10, font=BODY, color=FG_HIGH)
        add_text(s, f"{pct} %",
                 left=Inches(4.6), top=y + Inches(0.04), width=Inches(2.0), height=Inches(0.3),
                 size=10, bold=True, font=MONO, color=accent, align=PP_ALIGN.RIGHT)

    # Format mix — right card
    add_text(s, "Format mix (KPI #6 evidence)",
             left=Inches(7.1), top=Inches(3.3), width=Inches(5.55), height=Inches(0.3),
             size=12, bold=True, font=DISPLAY, color=FG_HIGH)
    formats = [("Reels", 4), ("Carousels", 4), ("Single image", 5), ("Stories", 2)]
    for i, (fname, count) in enumerate(formats):
        y = Inches(3.7 + i * 0.42)
        add_card(s, left=Inches(7.1), top=y, width=Inches(5.55), height=Inches(0.36),
                 accent=ACCENT_VIOLET)
        add_text(s, fname,
                 left=Inches(7.3), top=y + Inches(0.06), width=Inches(2.5), height=Inches(0.3),
                 size=10, font=BODY, color=FG_HIGH)
        # bar
        bw_max = Inches(2.0); bw_cur = Inches(bw_max.inches * (count / 5.0))
        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(9.85), y + Inches(0.1), bw_cur, Inches(0.16))
        bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_FUCHSIA
        bar.line.fill.background()
        add_text(s, str(count),
                 left=Inches(11.9), top=y + Inches(0.06), width=Inches(0.6), height=Inches(0.3),
                 size=11, bold=True, font=MONO, color=ACCENT_FUCHSIA, align=PP_ALIGN.RIGHT)
    add_text(s, "4 unique formats · target ≥ 3 ✅",
             left=Inches(7.1), top=Inches(5.4), width=Inches(5.55), height=Inches(0.3),
             size=10, bold=True, font=BODY, color=ACCENT_EMERALD)

    # Sample headlines from the calendar
    add_text(s, "Sample headlines from the 14-day calendar",
             left=Inches(0.7), top=Inches(5.85), width=Inches(12), height=Inches(0.3),
             size=12, bold=True, font=DISPLAY, color=FG_HIGH)
    samples = [
        ("D1 Mon · Reel",        "\"You hit play on a 90-min video…\" — 30 s teaser",      ACCENT_VIOLET),
        ("D5 Fri · Meme",        "\"Me: I'll watch it later. Video, 2 hrs later, still unwatched.\"", ACCENT_AMBER),
        ("D7 Sun · Carousel",    "\"5 video genres VidIQ is best at\" — lectures · podcasts · trading · medical · legal", ACCENT_CYAN),
        ("D11 Thu · Quote",      "\"I cut my lecture-review time from 2 hrs to 8 min.\" — Beta tester", ACCENT_EMERALD),
    ]
    for i, (when, text, accent) in enumerate(samples):
        y = Inches(6.2 + (i // 2) * 0.45)
        x = Inches(0.7 + (i % 2) * 6.0)
        add_card(s, left=x, top=y, width=Inches(5.95), height=Inches(0.4),
                 accent=accent)
        add_text(s, when,
                 left=x + Inches(0.15), top=y + Inches(0.08), width=Inches(1.65), height=Inches(0.3),
                 size=9, bold=True, font=MONO, color=accent)
        add_text(s, text,
                 left=x + Inches(1.85), top=y + Inches(0.08), width=Inches(4.0), height=Inches(0.3),
                 size=9, font=BODY, color=FG_HIGH)

    set_notes(s, "Aim: educate (E) → show product (P) → use case (U) — front-loaded, with weekly meme + "
                 "social proof to keep the tone human. 4 unique formats clears KPI #6 (≥3). All 28 posts "
                 "queued in Business Suite Planner with publish dates 1 year in future (showcase mode). "
                 "Time: 50 s.")


def slide_24_meta_ads_4_sets(prs):
    s = add_blank_slide(prs)
    add_header(s, "Pillar 2 · Meta Ads · 4 ad sets",
               "$20/day · $280 over 14 days · CBO · learn → optimise.",
               accent=ACCENT_FUCHSIA)
    add_footer(s, "Pillar 2 · Meta Ads · 04-meta-ads-plan.md")

    # Architecture banner
    add_card(s, left=Inches(0.7), top=Inches(2.35), width=Inches(11.95), height=Inches(0.7),
             accent=ACCENT_FUCHSIA)
    arch = "Campaign · Traffic (D1–D7) → Conversions \"StartedAnalysis\" (D8–D14)  ·  Bid: Highest Volume → CPR $1.50  ·  Attribution: 7-day click / 1-day view"
    add_text(s, arch,
             left=Inches(0.95), top=Inches(2.55), width=Inches(11.45), height=Inches(0.4),
             size=10, font=BODY, color=FG_HIGH, line_spacing=1.3)

    # 4 ad set cards (2×2)
    sets = [
        ("AS 1 · Students",       "🎓",
         "$5/day · $70 total",
         "PK · IN · BD · LK · EG · NG · PH",
         "18 – 26 · all genders · English",
         "Khan Academy · Coursera · edX · Notion · Quizlet · Anki + LAL 1 % off StartedAnalysis",
         "IG Reels + Feed · FB Reels + Feed",
         "Aim: cheapest CPM tier-2 markets · build LAL seed",
         ACCENT_VIOLET),
        ("AS 2 · Creators",       "🎥",
         "$5/day · $70 total",
         "US · UK · CA · AU · UAE · PK · IN",
         "22 – 40 · all genders",
         "YouTube creator economy · Descript · Riverside · OBS · DaVinci · MrBeast · Ali Abdaal",
         "IG Reels (primary) · IG Feed · FB Reels",
         "Aim: highest LTV proxy · mid-funnel research intent",
         ACCENT_FUCHSIA),
        ("AS 3 · Knowledge wkrs", "💼",
         "$5/day · $70 total",
         "US · UK · CA · DE · SG · UAE · PK",
         "25 – 45 · job titles: PM/Researcher/Analyst",
         "Notion · Linear · Slack · Loom · Otter.ai · Read.ai · MIT Tech Review",
         "FB Feed · IG Feed · FB Right Column · Messenger Inbox",
         "Aim: premium CPM markets · higher conversion value",
         ACCENT_CYAN),
        ("AS 4 · Retargeting",    "🔁",
         "$5/day · $70 total",
         "All retargetable geos",
         "Custom audiences only · 1 / day frequency cap",
         "30-d /analyze visitors · 90-d page engagers · 75 %+ video viewers",
         "All Meta surfaces",
         "Aim: 4× ROAS uplift vs cold (industry benchmark)",
         ACCENT_EMERALD),
    ]
    cw = Inches(5.95); ch = Inches(1.85); cg = Inches(0.05)
    sx = Inches(0.7); sy = Inches(3.2)
    for i, (title, emoji, budget, geo, age, interests, placements, aim, accent) in enumerate(sets):
        col = i % 2; row = i // 2
        x = sx + (cw + cg) * col
        y = sy + (ch + cg) * row
        add_card(s, left=x, top=y, width=cw, height=ch, accent=accent)
        # title bar
        add_text(s, f"{emoji}  {title}",
                 left=x + Inches(0.18), top=y + Inches(0.12),
                 width=Inches(3.3), height=Inches(0.3),
                 size=13, bold=True, font=DISPLAY, color=accent)
        add_text(s, budget,
                 left=x + Inches(3.5), top=y + Inches(0.16),
                 width=Inches(2.3), height=Inches(0.25),
                 size=9, bold=True, font=MONO, color=FG_MID, align=PP_ALIGN.RIGHT)
        # rows of facts
        rows = [("Geo",       geo),
                ("Audience",  age),
                ("Targeting", interests),
                ("Placements",placements)]
        for j, (k, v) in enumerate(rows):
            ry = y + Inches(0.5 + j * 0.27)
            add_text(s, k.upper(),
                     left=x + Inches(0.18), top=ry,
                     width=Inches(0.95), height=Inches(0.22),
                     size=7, bold=True, font=BODY, color=FG_LOW)
            add_text(s, v,
                     left=x + Inches(1.2), top=ry,
                     width=cw - Inches(1.4), height=Inches(0.27),
                     size=8, font=BODY, color=FG_HIGH, line_spacing=1.25)
        # aim — bottom band
        add_text(s, aim,
                 left=x + Inches(0.18), top=y + Inches(1.55),
                 width=cw - Inches(0.36), height=Inches(0.27),
                 size=8, bold=True, font=BODY, color=accent, line_spacing=1.2)

    # Footer KPI strip
    kpi = "Targets: CTR > 1.5 %  ·  CPC < $0.30  ·  CPA < $1.50  ·  Frequency < 2.5 (else swap creative)"
    add_text(s, kpi,
             left=Inches(0.7), top=Inches(7.05), width=Inches(12), height=Inches(0.3),
             size=9, font=BODY, color=FG_LOW)

    set_notes(s, "Aim per ad set is what the markers want to hear. Students = volume + LAL seed. "
                 "Creators = LTV. Knowledge workers = premium CPM. Retargeting = ROAS uplift. "
                 "Showcase mode means we built each to the *Review* page in Ads Manager and screenshotted "
                 "the Summary screen — that is the rubric deliverable. Time: 75 s.")


def slide_25_conversation_layer(prs):
    s = add_blank_slide(prs)
    add_header(s, "Pillar 2 · Always-on conversation",
               "Welcome · away · 7 FAQs · 5 saved replies — KPI #9 evidence.",
               accent=ACCENT_FUCHSIA)
    add_footer(s, "Pillar 2 · Inbox · 05-social-templates.md")

    # Three quote cards across
    quotes = [
        ("Greeting · Messenger + IG DM",
         "👋 Hi, {first_name} — welcome to VidIQ.",
         "We turn any YouTube video or live stream into a transcript, time-stamped summary, and a chat "
         "you can ground every answer in. Reply with a video link and we'll show you. Or try it now → "
         "vidiq-two.vercel.app/analyze",
         ACCENT_VIOLET),
        ("Away · outside hours",
         "Thanks 🌙 — we reply within 24 hours.",
         "Don't wait for us — paste any video URL into vidiq-two.vercel.app/analyze and you'll have a "
         "full analysis in under a minute. Free. No card.",
         ACCENT_FUCHSIA),
        ("FAQ #1 · price / free",
         "Trigger: \"price · pricing · cost · free\"",
         "VidIQ runs on free-tier APIs (Google Gemini + Groq) — your first analyses cost you nothing. "
         "We'll announce paid plans before any pricing change.",
         ACCENT_CYAN),
    ]
    qw = Inches(4.05); qh = Inches(2.3); qg = Inches(0.12)
    sx = Inches(0.7); sy = Inches(2.4)
    for i, (eyebrow, headline, body, accent) in enumerate(quotes):
        x = sx + (qw + qg) * i
        add_card(s, left=x, top=sy, width=qw, height=qh, accent=accent)
        add_text(s, eyebrow.upper(),
                 left=x + Inches(0.2), top=sy + Inches(0.15),
                 width=qw - Inches(0.4), height=Inches(0.25),
                 size=8, bold=True, font=BODY, color=accent)
        add_text(s, headline,
                 left=x + Inches(0.2), top=sy + Inches(0.45),
                 width=qw - Inches(0.4), height=Inches(0.6),
                 size=14, bold=True, font=DISPLAY, color=FG_HIGH, line_spacing=1.15)
        add_text(s, body,
                 left=x + Inches(0.2), top=sy + Inches(1.05),
                 width=qw - Inches(0.4), height=Inches(1.2),
                 size=9, font=BODY, color=FG_MID, line_spacing=1.35)

    # FAQ list — 6 more triggers
    add_text(s, "Other FAQs configured (4 more covered)",
             left=Inches(0.7), top=Inches(4.95), width=Inches(12), height=Inches(0.3),
             size=12, bold=True, font=DISPLAY, color=FG_HIGH)
    more_faqs = [
        ("how does it work",   "3 steps: paste URL · transcribe + analyse + summarise · read & chat with citations"),
        ("languages · urdu",   "YouTube transcript when present (100+ langs) · faster-whisper fallback"),
        ("live · stream",      "Yes — vidiq-two.vercel.app/live · rolling summaries on YouTube Live + webinars"),
        ("safe · privacy",     "Stored only in your account · we don't redistribute or train models on your content"),
    ]
    for i, (trigger, reply) in enumerate(more_faqs):
        col = i % 2; row = i // 2
        x = Inches(0.7 + col * 6.0)
        y = Inches(5.3 + row * 0.43)
        add_card(s, left=x, top=y, width=Inches(5.95), height=Inches(0.38),
                 accent=ACCENT_VIOLET)
        add_text(s, trigger,
                 left=x + Inches(0.18), top=y + Inches(0.07),
                 width=Inches(1.85), height=Inches(0.3),
                 size=9, bold=True, font=MONO, color=ACCENT_VIOLET)
        add_text(s, reply,
                 left=x + Inches(2.05), top=y + Inches(0.07),
                 width=Inches(3.8), height=Inches(0.3),
                 size=8, font=BODY, color=FG_HIGH)

    # Saved replies row
    add_text(s, "Saved replies for comments  ·  5 templates",
             left=Inches(0.7), top=Inches(6.25), width=Inches(12), height=Inches(0.3),
             size=12, bold=True, font=DISPLAY, color=FG_HIGH)
    saved = ["Compliment", "Curious \"how?\"", "Sceptical", "Tag-a-friend", "Feature request"]
    sw = Inches(2.42); sh = Inches(0.4); sg = Inches(0.05)
    for i, label in enumerate(saved):
        x = Inches(0.7) + (sw + sg) * i
        add_card(s, left=x, top=Inches(6.6), width=sw, height=sh, accent=ACCENT_EMERALD)
        add_text(s, label,
                 left=x, top=Inches(6.7), width=sw, height=Inches(0.25),
                 size=10, bold=True, font=BODY, color=ACCENT_EMERALD, align=PP_ALIGN.CENTER)

    set_notes(s, "Aim: KPI #9 — Automated message / Welcome note. Greeting + away + 7 FAQs + 5 saved "
                 "replies are all configured in Meta Business Suite → Inbox → Automations. Screenshots "
                 "in marketing/assets/screenshots/ (auto-greeting · auto-away · auto-faq · saved-replies). "
                 "Time: 45 s.")


def slide_26_keyword_seo(prs):
    s = add_blank_slide(prs)
    add_header(s, "Pillar 4 · Keyword + on-page SEO",
               "18 keywords · 5 head + 13 long-tail · 4 on-page elements audited.",
               accent=ACCENT_EMERALD)
    add_footer(s, "Pillar 4 · SEO · 06-keyword-research.md · 07-onpage-seo-report.md")

    # Left column — KW funnel 3 tiers
    add_text(s, "Keyword funnel — 18 KWs",
             left=Inches(0.7), top=Inches(2.35), width=Inches(6), height=Inches(0.35),
             size=14, bold=True, font=DISPLAY, color=FG_HIGH)

    tiers = [
        ("Head · KW 1-5",
         "Saturated → bid via Google Ads",
         [("ai video summarizer",    "27.1 k", 48),
          ("youtube video summarizer","22.2 k", 42),
          ("summarize youtube video", "18.1 k", 39),
          ("ai for video transcripts", "4.4 k", 31),
          ("live stream ai summary",   "1.6 k", 24)],
         ACCENT_VIOLET),
        ("Long-tail · KD < 15",
         "Rankable in 2-3 mo via on-page",
         [("summarize a 2-hr youtube video",            "2.4 k", 19),
          ("how to summarize a long youtube lecture",   "0.9 k", 12),
          ("study with youtube ai summary",             "0.7 k", 13),
          ("ai timestamp summary youtube",              "0.6 k", 11)],
         ACCENT_FUCHSIA),
        ("Niche · vertical-domain",
         "First-mover · zero competition",
         [("summarize trading livestream ai", "0.21 k", 8),
          ("medical lecture video summarizer","0.32 k", 14),
          ("webinar to text ai summary",      "1.1 k",  25)],
         ACCENT_CYAN),
    ]
    ty = Inches(2.78)
    for tier, hint, rows, accent in tiers:
        add_card(s, left=Inches(0.7), top=ty, width=Inches(6.0), height=Inches(0.35 + 0.27 * len(rows)),
                 accent=accent)
        add_text(s, tier,
                 left=Inches(0.85), top=ty + Inches(0.07), width=Inches(3.5), height=Inches(0.25),
                 size=11, bold=True, font=DISPLAY, color=accent)
        add_text(s, hint,
                 left=Inches(4.0), top=ty + Inches(0.09), width=Inches(2.6), height=Inches(0.25),
                 size=8, font=BODY, color=FG_LOW, align=PP_ALIGN.RIGHT)
        for j, (kw, msv, kd) in enumerate(rows):
            ry = ty + Inches(0.34 + j * 0.27)
            add_text(s, kw,
                     left=Inches(0.85), top=ry, width=Inches(3.6), height=Inches(0.25),
                     size=9, font=BODY, color=FG_HIGH)
            add_text(s, msv,
                     left=Inches(4.55), top=ry, width=Inches(0.85), height=Inches(0.25),
                     size=9, font=MONO, color=FG_MID, align=PP_ALIGN.RIGHT)
            add_text(s, f"KD {kd}",
                     left=Inches(5.5), top=ry, width=Inches(1.1), height=Inches(0.25),
                     size=9, bold=True, font=MONO, color=accent, align=PP_ALIGN.RIGHT)
        ty = ty + Inches(0.4 + 0.27 * len(rows))

    # Right column — On-page audit
    add_text(s, "On-page audit · 4 elements",
             left=Inches(7.1), top=Inches(2.35), width=Inches(5.5), height=Inches(0.35),
             size=14, bold=True, font=DISPLAY, color=FG_HIGH)

    audit = [
        ("Meta tags",      "✅",
         "%s · VidIQ title template · OG + Twitter + canonical on every page · JSON-LD: Organization + WebSite + SoftwareApplication"),
        ("Alt text",       "✅",
         "Every functional image has descriptive alt · decorative auroras `aria-hidden=\"true\"` · keyframe alts include timestamp"),
        ("Header tree",    "✅",
         "1 H1 per route · logical H2/H3 nesting · validated across /, /analyze, /live, /library, /videos/[id]"),
        ("KW placement",   "✅",
         "Primary KW + 3-5 secondary per page · cluster /analyze around KW 3, 6, 9, 10, 18 · /live around KW 5, 12, 13"),
    ]
    for i, (label, mark, body) in enumerate(audit):
        y = Inches(2.78 + i * 0.95)
        add_card(s, left=Inches(7.1), top=y, width=Inches(5.55), height=Inches(0.85),
                 accent=ACCENT_EMERALD)
        add_text(s, label,
                 left=Inches(7.3), top=y + Inches(0.1), width=Inches(2.5), height=Inches(0.3),
                 size=12, bold=True, font=DISPLAY, color=FG_HIGH)
        add_text(s, mark,
                 left=Inches(11.95), top=y + Inches(0.1), width=Inches(0.5), height=Inches(0.3),
                 size=14, bold=True, font=BODY, color=ACCENT_EMERALD, align=PP_ALIGN.RIGHT)
        add_text(s, body,
                 left=Inches(7.3), top=y + Inches(0.4), width=Inches(5.2), height=Inches(0.45),
                 size=8, font=BODY, color=FG_MID, line_spacing=1.35)

    # Footer rationale
    add_text(s, "Strategy: lead with long-tail KD<15 (rankable) + niche vertical-domain (first-mover). Skip head-term war with NoteGPT/Eightify — bid those via Google Ads instead.",
             left=Inches(0.7), top=Inches(6.95), width=Inches(12), height=Inches(0.5),
             size=9, font=BODY, color=FG_LOW, line_spacing=1.4)

    set_notes(s, "Aim: own long-tail + niche organically, bid head terms via paid. 18 KWs total — 5 short / "
                 "13 long-tail clears KPI #14 + #15. On-page audit covers all 4 rubric elements (meta, alt, "
                 "headers, KW). Verifiable in DevTools → Elements → <head> on the live URL. Time: 60 s.")


def slide_27_google_ads(prs):
    s = add_blank_slide(prs)
    add_header(s, "Pillar 4 · Google Ads · 3 campaigns",
               "$15/day · $210 over 14 days · Search + YouTube + Performance Max.",
               accent=ACCENT_EMERALD)
    add_footer(s, "Pillar 4 · Google Ads · 08-google-ads-plan.md")

    # 3 campaign cards
    campaigns = [
        ("Campaign 1", "📣  Search · Brand cluster",
         "$5 / day · $70 total",
         "Maximise Conv → Target CPA $1.50",
         "4 ad groups · 15 headlines + 4 desc / RSA",
         [
             ("Group A", "AI video summariser (head)",       "KW 1, 2, 3, 6, 8, 18"),
             ("Group B", "Chat / extract moments",            "KW 4, 9, 10"),
             ("Group C", "Live + webinar summary",            "KW 5, 12, 13"),
             ("Group D", "Students long-tail",                "KW 7, 14"),
         ],
         "Aim: lowest-funnel intent at cheapest CPC · Negative KW: vidiq.com (different brand)",
         ACCENT_VIOLET),
        ("Campaign 2", "🎬  YouTube · In-Stream",
         "$4 / day · $56 total",
         "Target CPM $4 (awareness)",
         "30 s skippable + 6 s bumper",
         [
             ("Audience",  "Custom Intent (KW 1, 2, 3, 6, 7, 14) + In-Market: Education software"),
             ("Topics",    "Online Education · Higher Education · YouTube"),
             ("Whitelist", "Khan Academy · MIT OCW · Lex Fridman · MKBHD · Ali Abdaal"),
         ],
         "Aim: brand recall + supports search-intent layer",
         ACCENT_FUCHSIA),
        ("Campaign 3", "🌐  Performance Max (gated)",
         "$6 / day · $84 (starts ≈ D8)",
         "Maximise Conversion Value",
         "3 asset groups · Final URL expansion OFF",
         [
             ("Group 1", "Students · 5 headlines / 5 desc / 5 images / 1 video"),
             ("Group 2", "Creators · same structure"),
             ("Group 3", "Knowledge workers · same structure"),
         ],
         "Aim: scale once 30 conversions banked / 7 d · automation finds incremental reach",
         ACCENT_CYAN),
    ]
    cw = Inches(3.93); ch = Inches(4.05); cg = Inches(0.08)
    sx = Inches(0.7); sy = Inches(2.35)
    for i, (label, title, budget, bid, creative, rows, aim, accent) in enumerate(campaigns):
        x = sx + (cw + cg) * i
        add_card(s, left=x, top=sy, width=cw, height=ch, accent=accent)
        add_text(s, label.upper(),
                 left=x + Inches(0.18), top=sy + Inches(0.13),
                 width=cw - Inches(0.36), height=Inches(0.25),
                 size=8, bold=True, font=BODY, color=accent)
        add_text(s, title,
                 left=x + Inches(0.18), top=sy + Inches(0.4),
                 width=cw - Inches(0.36), height=Inches(0.45),
                 size=14, bold=True, font=DISPLAY, color=FG_HIGH, line_spacing=1.15)
        # spec rows
        specs = [("Budget",    budget),
                 ("Bid",       bid),
                 ("Creative",  creative)]
        for j, (k, v) in enumerate(specs):
            ry = sy + Inches(0.95 + j * 0.34)
            add_text(s, k.upper(),
                     left=x + Inches(0.18), top=ry,
                     width=Inches(0.85), height=Inches(0.22),
                     size=7, bold=True, font=BODY, color=FG_LOW)
            add_text(s, v,
                     left=x + Inches(1.08), top=ry,
                     width=cw - Inches(1.26), height=Inches(0.32),
                     size=8, font=BODY, color=FG_HIGH, line_spacing=1.25)
        # Inner card with ad-group rows
        inner_y = sy + Inches(2.05)
        inner_h = Inches(1.5)
        inner = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    x + Inches(0.15), inner_y,
                                    cw - Inches(0.3), inner_h)
        inner.fill.solid(); inner.fill.fore_color.rgb = RGBColor(0x0F, 0x0A, 0x1C)
        inner.line.color.rgb = accent; inner.line.width = Pt(0.4)
        inner.adjustments[0] = 0.05
        for j, row in enumerate(rows):
            ry = inner_y + Inches(0.13 + j * 0.34)
            if len(row) == 3:
                k, v1, v2 = row
                add_text(s, k,
                         left=x + Inches(0.3), top=ry,
                         width=Inches(0.85), height=Inches(0.22),
                         size=8, bold=True, font=MONO, color=accent)
                add_text(s, v1,
                         left=x + Inches(1.15), top=ry,
                         width=cw - Inches(1.55), height=Inches(0.22),
                         size=8, font=BODY, color=FG_HIGH)
                add_text(s, v2,
                         left=x + Inches(0.3), top=ry + Inches(0.16),
                         width=cw - Inches(0.6), height=Inches(0.22),
                         size=7, font=MONO, color=FG_LOW)
            else:
                k, v = row
                add_text(s, k,
                         left=x + Inches(0.3), top=ry,
                         width=Inches(1.0), height=Inches(0.22),
                         size=8, bold=True, font=MONO, color=accent)
                add_text(s, v,
                         left=x + Inches(1.3), top=ry,
                         width=cw - Inches(1.55), height=Inches(0.32),
                         size=7, font=BODY, color=FG_HIGH, line_spacing=1.2)
        # Aim
        add_text(s, aim,
                 left=x + Inches(0.18), top=sy + Inches(3.65),
                 width=cw - Inches(0.36), height=Inches(0.4),
                 size=8, bold=True, font=BODY, color=accent, line_spacing=1.3)

    # Footer KPI strip
    kpi = "Targets: Search CTR > 6 %  ·  Search CPC < $0.40  ·  Conv rate > 4 %  ·  YouTube view-rate > 25 %  ·  YouTube CPV < $0.020"
    add_text(s, kpi,
             left=Inches(0.7), top=Inches(6.7), width=Inches(12), height=Inches(0.3),
             size=9, font=BODY, color=FG_LOW)

    set_notes(s, "Aim per campaign: Search = lowest-funnel conversion volume. YouTube = awareness layer. "
                 "PMax = scale once we have signal (30 conv / 7 rolling days). PMax is gated — until then, "
                 "spend goes into Search. UTM strategy lets GA4 stitch every click → conversion. Time: 75 s.")


def slide_28_budget(prs):
    s = add_blank_slide(prs)
    add_header(s, "Pillar 5 · Budget distribution",
               "₨ 250 000 PKR · ~$890 USD · 14-day flight · ₨ 0 actual (showcase).",
               accent=ACCENT_AMBER)
    add_footer(s, "Pillar 5 · Budget · 09-budget.md · PKR-primary")

    # Top-line line items (left)
    add_text(s, "Top-line (₨ 250,000)",
             left=Inches(0.7), top=Inches(2.35), width=Inches(5.0), height=Inches(0.3),
             size=12, bold=True, font=DISPLAY, color=FG_HIGH)
    items = [
        ("Meta Ads",            "₨ 95,000",  38, ACCENT_FUCHSIA),
        ("Google Ads",          "₨ 70,000",  28, ACCENT_EMERALD),
        ("Creative production", "₨ 40,000",  16, ACCENT_VIOLET),
        ("Tools & SaaS",        "₨ 17,500",   7, ACCENT_CYAN),
        ("Influencer seeding",  "₨ 20,000",   8, ACCENT_AMBER),
        ("Contingency 10 %",    "₨ 7,500",    3, FG_MID),
    ]
    for i, (name, pkr, pct, accent) in enumerate(items):
        y = Inches(2.7 + i * 0.42)
        add_card(s, left=Inches(0.7), top=y, width=Inches(5.0), height=Inches(0.36),
                 accent=accent)
        add_text(s, name,
                 left=Inches(0.85), top=y + Inches(0.06), width=Inches(2.0), height=Inches(0.3),
                 size=10, font=BODY, color=FG_HIGH)
        # bar
        bw_max = Inches(1.2); bw_cur = Inches(bw_max.inches * (pct / 38.0))
        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(2.95), y + Inches(0.1), bw_cur, Inches(0.16))
        bar.fill.solid(); bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
        add_text(s, pkr,
                 left=Inches(4.2), top=y + Inches(0.06), width=Inches(1.15), height=Inches(0.3),
                 size=10, bold=True, font=MONO, color=accent, align=PP_ALIGN.RIGHT)
        add_text(s, f"{pct}%",
                 left=Inches(5.32), top=y + Inches(0.06), width=Inches(0.4), height=Inches(0.3),
                 size=9, font=MONO, color=FG_LOW, align=PP_ALIGN.RIGHT)

    # Paid media split (centre)
    add_text(s, "Paid media split",
             left=Inches(6.0), top=Inches(2.35), width=Inches(3.4), height=Inches(0.3),
             size=12, bold=True, font=DISPLAY, color=FG_HIGH)
    # Meta ad sets
    add_text(s, "META · ₨ 95,000",
             left=Inches(6.0), top=Inches(2.7), width=Inches(3.4), height=Inches(0.3),
             size=8, bold=True, font=BODY, color=ACCENT_FUCHSIA)
    meta_split = [("Students ₨23,750", 25), ("Creators ₨23,750", 25), ("Knowledge ₨23,750", 25), ("Retarget ₨23,750", 25)]
    bx = Inches(6.0); by = Inches(3.0); bw = Inches(3.4); bh = Inches(0.4)
    cum = 0
    meta_colors = [ACCENT_VIOLET, ACCENT_FUCHSIA, ACCENT_CYAN, ACCENT_EMERALD]
    for j, (label, p) in enumerate(meta_split):
        seg_w = Inches(bw.inches * p / 100.0)
        seg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(bx.inches + bw.inches * cum / 100.0), by,
                                  seg_w, bh)
        seg.fill.solid(); seg.fill.fore_color.rgb = meta_colors[j]
        seg.line.fill.background()
        cum += p
    # Meta legend
    for j, (label, p) in enumerate(meta_split):
        ly = Inches(3.5 + j * 0.27)
        chip = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.05), ly + Inches(0.07),
                                   Inches(0.12), Inches(0.12))
        chip.fill.solid(); chip.fill.fore_color.rgb = meta_colors[j]
        chip.line.fill.background()
        add_text(s, label,
                 left=Inches(6.25), top=ly + Inches(0.04), width=Inches(3.2), height=Inches(0.25),
                 size=9, font=BODY, color=FG_HIGH)

    # Google split
    add_text(s, "GOOGLE · ₨ 70,000",
             left=Inches(6.0), top=Inches(4.85), width=Inches(3.4), height=Inches(0.3),
             size=8, bold=True, font=BODY, color=ACCENT_EMERALD)
    g_split = [("Search ₨23,100", 33), ("YouTube ₨18,900", 27), ("PMax ₨28,000", 40)]
    by2 = Inches(5.15)
    cum = 0
    g_colors = [ACCENT_VIOLET, ACCENT_FUCHSIA, ACCENT_CYAN]
    for j, (label, p) in enumerate(g_split):
        seg_w = Inches(bw.inches * p / 100.0)
        seg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(bx.inches + bw.inches * cum / 100.0), by2,
                                  seg_w, bh)
        seg.fill.solid(); seg.fill.fore_color.rgb = g_colors[j]
        seg.line.fill.background()
        cum += p
    for j, (label, p) in enumerate(g_split):
        ly = Inches(5.65 + j * 0.27)
        chip = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.05), ly + Inches(0.07),
                                   Inches(0.12), Inches(0.12))
        chip.fill.solid(); chip.fill.fore_color.rgb = g_colors[j]
        chip.line.fill.background()
        add_text(s, label,
                 left=Inches(6.25), top=ly + Inches(0.04), width=Inches(3.2), height=Inches(0.25),
                 size=9, font=BODY, color=FG_HIGH)

    # Phasing timeline (right)
    add_text(s, "Spend phasing",
             left=Inches(9.7), top=Inches(2.35), width=Inches(3.0), height=Inches(0.3),
             size=12, bold=True, font=DISPLAY, color=FG_HIGH)

    phases = [
        ("D1 – D7",   "Learn",
         "Highest Volume bid · LP-views objective · cheapest data while pixel learns",
         ACCENT_VIOLET),
        ("D5",        "Creative swap fund",
         "10 % of Meta budget reserved for fatigue-rotation creative",
         ACCENT_AMBER),
        ("D8 – D14",  "Optimise",
         "Switch to Conversions + CPR $1.50 cap · PMax unlocks if 30 conv banked",
         ACCENT_FUCHSIA),
        ("D14",       "Wrap",
         "Negative-KW sweep · screenshot all reports · publish learnings",
         ACCENT_EMERALD),
    ]
    for i, (when, label, body, accent) in enumerate(phases):
        y = Inches(2.7 + i * 1.05)
        add_card(s, left=Inches(9.7), top=y, width=Inches(3.0), height=Inches(0.95),
                 accent=accent)
        add_text(s, when,
                 left=Inches(9.85), top=y + Inches(0.1), width=Inches(2.7), height=Inches(0.25),
                 size=9, bold=True, font=MONO, color=accent)
        add_text(s, label,
                 left=Inches(9.85), top=y + Inches(0.32), width=Inches(2.7), height=Inches(0.25),
                 size=12, bold=True, font=DISPLAY, color=FG_HIGH)
        add_text(s, body,
                 left=Inches(9.85), top=y + Inches(0.55), width=Inches(2.7), height=Inches(0.4),
                 size=8, font=BODY, color=FG_MID, line_spacing=1.3)

    # Showcase note footer + PKR rationale
    add_text(s, "Why ₨ 250 K? Pakistani SaaS launch budgets cluster at ₨ 150-400 K/month. PK Meta CPM ≈ ₨ 80-250 vs $5-15 in West (5-7× more reach per rupee). PK Google CPC ≈ ₨ 20-80 vs $1-3 USD.",
             left=Inches(0.7), top=Inches(6.85), width=Inches(12), height=Inches(0.3),
             size=8, font=BODY, color=ACCENT_AMBER, line_spacing=1.3)
    add_text(s, "Showcase note: actual spend = ₨ 0. Campaigns built to *Review* in Ads Manager and screenshotted — that is the brief deliverable, not paid reach. FX: 1 USD ≈ 280 PKR.",
             left=Inches(0.7), top=Inches(7.15), width=Inches(12), height=Inches(0.3),
             size=8, font=BODY, color=FG_LOW, line_spacing=1.3)

    set_notes(s, "Top-line ₨ 250 000 (~$890) splits 66 % paid media (Meta + Google), 16 % creative, 7 % tools, "
                 "8 % influencer, 3 % contingency. We bill in PKR because we run from Pakistan; PK CPM/CPC is "
                 "5-7× cheaper than Western markets, so this PKR budget converts to Western-equivalent reach "
                 "of ~$3-4K spend. Spend phasing matches Meta's objective ladder (LP-views D1-7 → Conversions "
                 "D8-14) and Google's PMax gating (30 conv / 7 d). Time: 60 s.")


def slide_29_measurement(prs):
    s = add_blank_slide(prs)
    add_header(s, "Conversion tracking & measurement",
               "Meta Pixel + CAPI · Google gtag · GA4 · same event names across stack.",
               accent=ACCENT_CYAN)
    add_footer(s, "Pillars 2 + 4 · Tracking architecture")

    # Two columns — Meta + Google
    cards = [
        ("Meta · Pixel + CAPI",
         "Browser Pixel + server-side CAPI for iOS 14+ resilience",
         [
             ("PageView",                 "every page mounts"),
             ("ViewContent",              "/analyze page mounts · content_name='analyze_form'"),
             ("Lead",                     "'Start analysis' click — fired *before* mutation"),
             ("StartedAnalysis (custom)", "mutation 200 OK — Pixel + CAPI server-side from FastAPI"),
             ("CompleteRegistration",     "future — when auth ships"),
         ],
         ACCENT_FUCHSIA),
        ("Google · gtag → GA4",
         "Same event names · UTM-stitched in GA4",
         [
             ("started_analysis", "fires from /analyze onSuccess — same hook as Meta CAPI"),
             ("utm_source",       "google · meta · organic · email"),
             ("utm_medium",       "cpc · organic · social"),
             ("utm_campaign",     "search_brand · youtube_awareness · pmax_q1 · launch_wk1"),
             ("utm_term",         "auto-replaced by Google with the matched keyword"),
         ],
         ACCENT_EMERALD),
    ]
    cw = Inches(6.0); ch = Inches(3.0); cg = Inches(0.15)
    sx = Inches(0.7); sy = Inches(2.35)
    for i, (title, sub, rows, accent) in enumerate(cards):
        x = sx + (cw + cg) * i
        add_card(s, left=x, top=sy, width=cw, height=ch, accent=accent)
        add_text(s, title,
                 left=x + Inches(0.2), top=sy + Inches(0.13),
                 width=cw - Inches(0.4), height=Inches(0.4),
                 size=15, bold=True, font=DISPLAY, color=accent)
        add_text(s, sub,
                 left=x + Inches(0.2), top=sy + Inches(0.55),
                 width=cw - Inches(0.4), height=Inches(0.3),
                 size=9, font=BODY, color=FG_MID)
        for j, (k, v) in enumerate(rows):
            ry = sy + Inches(0.95 + j * 0.4)
            add_text(s, k,
                     left=x + Inches(0.2), top=ry,
                     width=Inches(2.2), height=Inches(0.3),
                     size=9, bold=True, font=MONO, color=accent)
            add_text(s, v,
                     left=x + Inches(2.45), top=ry,
                     width=cw - Inches(2.65), height=Inches(0.4),
                     size=8, font=BODY, color=FG_HIGH, line_spacing=1.3)

    # KPI targets row — combined targets
    add_text(s, "Cross-stack KPI targets — first 14 days",
             left=Inches(0.7), top=Inches(5.55), width=Inches(12), height=Inches(0.3),
             size=12, bold=True, font=DISPLAY, color=FG_HIGH)
    targets = [
        ("Meta · CTR",          "> 1.5 %", ACCENT_FUCHSIA),
        ("Meta · CPC",          "< $0.30", ACCENT_FUCHSIA),
        ("Meta · CPA",          "< $1.50", ACCENT_FUCHSIA),
        ("Google · Search CTR", "> 6 %",   ACCENT_EMERALD),
        ("Google · CPC",        "< $0.40", ACCENT_EMERALD),
        ("Google · Conv rate",  "> 4 %",   ACCENT_EMERALD),
    ]
    tw = Inches(2.0); th = Inches(0.7); tg = Inches(0.05)
    for i, (k, v, accent) in enumerate(targets):
        x = Inches(0.7) + (tw + tg) * i
        add_card(s, left=x, top=Inches(5.9), width=tw, height=th, accent=accent)
        add_text(s, k.upper(),
                 left=x + Inches(0.1), top=Inches(5.97), width=tw - Inches(0.2), height=Inches(0.28),
                 size=8, bold=True, font=BODY, color=accent, align=PP_ALIGN.CENTER)
        add_text(s, v,
                 left=x, top=Inches(6.25), width=tw, height=Inches(0.4),
                 size=15, bold=True, font=DISPLAY, color=FG_HIGH, align=PP_ALIGN.CENTER)

    # Showcase note footer
    add_text(s, "Showcase note: Pixel + CAPI not installed (no live spend = no events to capture). Plan presented as 'how we'd measure' — gate behind NEXT_PUBLIC_META_PIXEL_ID env var when going live.",
             left=Inches(0.7), top=Inches(6.85), width=Inches(12), height=Inches(0.4),
             size=9, font=BODY, color=FG_LOW, line_spacing=1.4)

    set_notes(s, "Measurement layer is identical across Meta and Google — same event name "
                 "(`StartedAnalysis` / `started_analysis`) fires from one place: the /analyze "
                 "useMutation onSuccess. Server-side CAPI gives iOS 14+ resilience the Pixel alone "
                 "doesn't. Time: 50 s.")


def slide_30_moats_roadmap(prs):
    s = add_blank_slide(prs)
    add_header(s, "Strategic moats + execution roadmap",
               "What we win on  ·  what we shipped  ·  what's next.",
               accent=ACCENT_VIOLET)
    add_footer(s, "Strategy · Execution · 10-competitive · 16-execution-playbook")

    # Left — 5 moats
    add_text(s, "Strategic moats",
             left=Inches(0.7), top=Inches(2.35), width=Inches(6), height=Inches(0.3),
             size=14, bold=True, font=DISPLAY, color=FG_HIGH)
    moats = [
        ("Multimodal grounding",
         "Vision + audio + LLM with timestamp citations. Neither NoteGPT nor Eightify ship this.",
         ACCENT_VIOLET),
        ("Live-stream pipeline",
         "Rolling summaries on YouTube Live, webinars and lectures. Owns KW cluster #5, #12, #13.",
         ACCENT_FUCHSIA),
        ("Vertical-domain modes",
         "Medical, legal, trading prompts adapt the analysis. First-mover — neither competitor offers it.",
         ACCENT_CYAN),
        ("Provider-agnostic backend",
         "Auto-failover Gemini ↔ Groq. Free-tier resilient. Switch providers via env var.",
         ACCENT_EMERALD),
        ("Citation-grounded chat",
         "Every claim cites the exact transcript moment + linked frame. YouTube's auto-summary is generic.",
         ACCENT_AMBER),
    ]
    for i, (title, body, accent) in enumerate(moats):
        y = Inches(2.7 + i * 0.86)
        add_card(s, left=Inches(0.7), top=y, width=Inches(6.0), height=Inches(0.78),
                 accent=accent)
        # number badge
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(0.85), y + Inches(0.18),
                                    Inches(0.4), Inches(0.4))
        badge.fill.solid(); badge.fill.fore_color.rgb = accent
        badge.line.fill.background()
        add_text(s, str(i + 1),
                 left=Inches(0.85), top=y + Inches(0.21),
                 width=Inches(0.4), height=Inches(0.35),
                 size=14, bold=True, font=DISPLAY, color=BG_DEEP, align=PP_ALIGN.CENTER)
        add_text(s, title,
                 left=Inches(1.4), top=y + Inches(0.1),
                 width=Inches(5.1), height=Inches(0.3),
                 size=12, bold=True, font=DISPLAY, color=accent)
        add_text(s, body,
                 left=Inches(1.4), top=y + Inches(0.4),
                 width=Inches(5.1), height=Inches(0.45),
                 size=8, font=BODY, color=FG_HIGH, line_spacing=1.3)

    # Right — 8-phase execution roadmap
    add_text(s, "Execution playbook · 8 phases",
             left=Inches(7.1), top=Inches(2.35), width=Inches(5.5), height=Inches(0.3),
             size=14, bold=True, font=DISPLAY, color=FG_HIGH)
    phases = [
        ("Phase 1", "Facebook Page created · cover · CTA",          "✅"),
        ("Phase 2", "Instagram Business linked",                     "✅"),
        ("Phase 3", "Welcome · away · 7 FAQs · 5 saved replies",     "✅"),
        ("Phase 4", "Adobe Firefly — 8 brand images IMG-1…IMG-8",    "✅"),
        ("Phase 5", "28 posts queued in Business Suite Planner",     "✅"),
        ("Phase 6", "Meta Ads · 4 ad sets to *Review* (drafts)",     "✅"),
        ("Phase 7", "Google Ads · 3 campaigns to *Review*",          "✅"),
        ("Phase 8", "Lighthouse + Search Console + final ZIP",       "⏳"),
    ]
    for i, (label, body, mark) in enumerate(phases):
        y = Inches(2.7 + i * 0.55)
        accent = ACCENT_EMERALD if mark == "✅" else ACCENT_AMBER
        add_card(s, left=Inches(7.1), top=y, width=Inches(5.55), height=Inches(0.48),
                 accent=accent)
        add_text(s, label,
                 left=Inches(7.25), top=y + Inches(0.13),
                 width=Inches(0.95), height=Inches(0.25),
                 size=9, bold=True, font=MONO, color=accent)
        add_text(s, body,
                 left=Inches(8.25), top=y + Inches(0.13),
                 width=Inches(3.85), height=Inches(0.3),
                 size=9, font=BODY, color=FG_HIGH)
        add_text(s, mark,
                 left=Inches(12.15), top=y + Inches(0.1),
                 width=Inches(0.4), height=Inches(0.3),
                 size=14, bold=True, font=BODY, color=accent, align=PP_ALIGN.CENTER)

    # Footer URLs strip
    urls = "Live · vidiq-two.vercel.app   ·   API · noumanhafeez11-vidiq-backend.hf.space   ·   Repo · github.com/noumanic/VidIQ   ·   KPI tracker · /marketing"
    add_text(s, urls,
             left=Inches(0.7), top=Inches(7.05), width=Inches(12), height=Inches(0.3),
             size=8, font=MONO, color=FG_LOW, align=PP_ALIGN.CENTER)

    set_notes(s, "Close on the moats + roadmap. The 5 moats are durable — citation grounding and "
                 "live-stream support are not commodity. 7 / 8 phases shipped — only Lighthouse audit + "
                 "final ZIP packaging remains. End on the live URL: vidiq-two.vercel.app. Time: 50 s.")


def slide_31_seo_evidence(prs):
    s = add_blank_slide(prs)
    add_header(s, "Pillar 4 · SEO evidence trail",
               "GSC verified · sitemap submitted · 3 URLs indexed · Lighthouse audited.",
               accent=ACCENT_EMERALD)
    add_footer(s, "Pillar 4 · SEO evidence")

    # 2x3 grid — 5 screenshots + 1 stats card
    cells = [
        # Row 1
        {"kind": "img", "file": "gsc-property-verified.png",
         "label": "Search Console · property verified",
         "caption": "HTML-tag method · vidiq-two.vercel.app",
         "accent": ACCENT_EMERALD},
        {"kind": "img", "file": "gsc-sitemap-success.png",
         "label": "Sitemap · submitted to GSC",
         "caption": "/sitemap.xml · 4 URLs discovered",
         "accent": ACCENT_VIOLET},
        {"kind": "img", "file": "gsc-url-inspection-home.png",
         "label": "URL inspection · indexing requested",
         "caption": "/, /analyze, /live · priority crawl queue",
         "accent": ACCENT_FUCHSIA},
        # Row 2
        {"kind": "img", "file": "pagespeed-mobile.png",
         "label": "PageSpeed · mobile",
         "caption": "Lighthouse audit · 4 categories",
         "accent": ACCENT_CYAN},
        {"kind": "img", "file": "pagespeed-desktop.png",
         "label": "PageSpeed · desktop",
         "caption": "Lighthouse audit · 4 categories",
         "accent": ACCENT_CYAN},
        # Stats card (no image)
        {"kind": "stats",
         "label": "Evidence summary",
         "rows": [
             ("Property",   "verified"),
             ("Sitemap",    "submitted · 4 URLs"),
             ("Inspector",  "/, /analyze, /live"),
             ("Lighthouse", "mobile + desktop"),
             ("Robots",     "vidiq-two.vercel.app"),
             ("JSON-LD",    "Org · WebSite · App"),
         ],
         "accent": ACCENT_AMBER},
    ]

    cw = Inches(3.95); ch = Inches(2.25); cg_x = Inches(0.18); cg_y = Inches(0.18)
    sx = Inches(0.7); sy = Inches(2.35)

    for i, cell in enumerate(cells):
        col = i % 3; row = i // 3
        x = sx + (cw + cg_x) * col
        y = sy + (ch + cg_y) * row
        accent = cell["accent"]

        add_card(s, left=x, top=y, width=cw, height=ch, accent=accent)

        # Label strip on top
        add_text(s, cell["label"].upper(),
                 left=x + Inches(0.18), top=y + Inches(0.13),
                 width=cw - Inches(0.36), height=Inches(0.25),
                 size=8, bold=True, font=BODY, color=accent)

        if cell["kind"] == "img":
            # Image centred in lower portion of card
            img_y = y + Inches(0.45)
            img_h = Inches(1.4)
            img_w = cw - Inches(0.4)
            add_image(s, asset(cell["file"]),
                      left=x + Inches(0.2), top=img_y,
                      width=img_w, height=img_h,
                      border_color=accent, border_pt=0.6)
            # Caption at bottom
            add_text(s, cell["caption"],
                     left=x + Inches(0.18), top=y + Inches(1.93),
                     width=cw - Inches(0.36), height=Inches(0.25),
                     size=8, font=BODY, color=FG_MID,
                     align=PP_ALIGN.CENTER)
        else:
            # Stats rows
            for j, (k, v) in enumerate(cell["rows"]):
                ry = y + Inches(0.5 + j * 0.27)
                add_text(s, k.upper(),
                         left=x + Inches(0.2), top=ry,
                         width=Inches(1.3), height=Inches(0.25),
                         size=8, bold=True, font=BODY, color=accent)
                add_text(s, v,
                         left=x + Inches(1.55), top=ry,
                         width=cw - Inches(1.75), height=Inches(0.25),
                         size=9, font=MONO, color=FG_HIGH)

    # Footer note
    add_text(s, "All evidence captured 2026-05-06 · live URL https://vidiq-two.vercel.app · screenshots in marketing/assets/seo-evidence/",
             left=Inches(0.7), top=Inches(7.05), width=Inches(12), height=Inches(0.3),
             size=9, font=BODY, color=FG_LOW, align=PP_ALIGN.CENTER)

    set_notes(s, "Pillar 4 evidence in one slide. GSC property verified via HTML-tag method "
                 "(NEXT_PUBLIC_GSC_VERIFICATION env var on Vercel). Sitemap submitted (4 URLs · "
                 "/, /analyze, /live, /library). Three priority URLs requested for indexing. "
                 "Lighthouse audited on mobile + desktop. Time: 45 s.")


def slide_32_logo_system(prs):
    s = add_blank_slide(prs)
    add_header(s, "Pillar 1 · Logo system",
               "Three lockups · scalable from 24 px favicon to print · brand-locked.",
               accent=ACCENT_VIOLET)
    add_footer(s, "Pillar 1 · Logo system")

    # ── Three large logo cards in a row ──────────────────────────────
    # Each card has its own contrast background so the logo is legible.
    logos = [
        {
            "file":   LOGO_MARK_DARK,
            "label":  "Mark · dark",
            "tagline":"For dark surfaces",
            "uses":   "Favicon · app icon · social avatar · slide footer · video end-card on dark backdrop",
            "bg":     BG_DEEP,
            "border": ACCENT_VIOLET,
            "is_text": False,
        },
        {
            "file":   LOGO_MARK_LIGHT,
            "label":  "Mark · light",
            "tagline":"For light surfaces",
            "uses":   "Print · presentation export · OG image (1200×630) · invoice + email signature",
            "bg":     RGBColor(0xF5, 0xEF, 0xFF),
            "border": ACCENT_FUCHSIA,
            "is_text": False,
        },
        {
            "file":   LOGO_TEXT,
            "label":  "Wordmark · horizontal lockup",
            "tagline":"For long-form layouts",
            "uses":   "Ad creatives · video end-cards · footer of long documents · website hero",
            "bg":     BG_CARD,
            "border": ACCENT_CYAN,
            "is_text": True,
        },
    ]
    cw = Inches(4.05); ch = Inches(3.7); cg = Inches(0.12)
    sx = Inches(0.7); sy = Inches(2.35)

    for i, item in enumerate(logos):
        x = sx + (cw + cg) * i

        # Outer accent border
        outer = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, sy, cw, ch)
        outer.fill.solid(); outer.fill.fore_color.rgb = BG_CARD
        outer.line.color.rgb = item["border"]; outer.line.width = Pt(1.5)
        outer.adjustments[0] = 0.04

        # Inner contrast plate for the logo (so light marks pop on light bg, dark on dark)
        inner_top = sy + Inches(0.3)
        inner_h = Inches(2.2)
        inner = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    x + Inches(0.25), inner_top,
                                    cw - Inches(0.5), inner_h)
        inner.fill.solid(); inner.fill.fore_color.rgb = item["bg"]
        inner.line.color.rgb = item["border"]; inner.line.width = Pt(0.5)
        inner.adjustments[0] = 0.05

        # Logo image — sized large enough to actually read
        if asset(item["file"]):
            if item["is_text"]:
                # Wordmark: spans card width with ~3:1 aspect
                logo_w = cw - Inches(0.9)
                logo_h = Inches(0.95)
                logo_x = x + (cw - logo_w) / 2
                logo_y = inner_top + (inner_h - logo_h) / 2
            else:
                # Square mark: 1:1 large
                logo_w = Inches(1.7)
                logo_h = Inches(1.7)
                logo_x = x + (cw - logo_w) / 2
                logo_y = inner_top + (inner_h - logo_h) / 2
            add_image(s, asset(item["file"]),
                      left=logo_x, top=logo_y,
                      width=logo_w, height=logo_h)

        # Label
        add_text(s, item["label"],
                 left=x + Inches(0.18), top=sy + Inches(2.65),
                 width=cw - Inches(0.36), height=Inches(0.35),
                 size=14, bold=True, font=DISPLAY, color=item["border"],
                 align=PP_ALIGN.CENTER)
        # Tagline
        add_text(s, item["tagline"].upper(),
                 left=x + Inches(0.18), top=sy + Inches(3.0),
                 width=cw - Inches(0.36), height=Inches(0.25),
                 size=8, bold=True, font=BODY, color=FG_LOW,
                 align=PP_ALIGN.CENTER)
        # Uses
        add_text(s, item["uses"],
                 left=x + Inches(0.22), top=sy + Inches(3.27),
                 width=cw - Inches(0.44), height=Inches(0.42),
                 size=8, font=BODY, color=FG_MID,
                 align=PP_ALIGN.CENTER, line_spacing=1.3)

    # ── Brand rules strip ────────────────────────────────────────────
    rules = [
        ("Clear-space",   "≥ 1× mark height on all four sides"),
        ("Minimum sizes", "24 px favicon · 40 px UI · 96 px print"),
        ("Colour rule",   "Never re-tint · use supplied PNG / SVG only"),
        ("File formats",  "SVG primary · PNG fallback · always 1:1 aspect"),
    ]
    rw = Inches(2.95); rh = Inches(0.7); rg = Inches(0.07)
    rsx = Inches(0.7); rsy = Inches(6.25)
    for i, (k, v) in enumerate(rules):
        x = rsx + (rw + rg) * i
        add_card(s, left=x, top=rsy, width=rw, height=rh, accent=ACCENT_AMBER)
        add_text(s, k.upper(),
                 left=x + Inches(0.12), top=rsy + Inches(0.08),
                 width=rw - Inches(0.24), height=Inches(0.22),
                 size=8, bold=True, font=BODY, color=ACCENT_AMBER)
        add_text(s, v,
                 left=x + Inches(0.12), top=rsy + Inches(0.32),
                 width=rw - Inches(0.24), height=Inches(0.32),
                 size=9, font=BODY, color=FG_HIGH, line_spacing=1.25)

    set_notes(s, "Three lockups: dark mark for dark surfaces (favicon, slide footers, social avatar). "
                 "Light mark for light surfaces (OG card, print). Wordmark for long-form (footer of docs, "
                 "ad creatives, video end-cards). Brand rules at the bottom: clear-space, minimum sizes, "
                 "never re-colour. SVG is primary; PNG fallback for legacy targets. Time: 50 s.")


def slide_33_logo_in_the_wild(prs):
    s = add_blank_slide(prs)
    add_header(s, "Pillar 1 · Logo in the wild",
               "Brand-mark applied across 7 surfaces — verified consistency.",
               accent=ACCENT_VIOLET)
    add_footer(s, "Pillar 1 · Logo applied")

    # ── Two columns: digital surfaces (left) + screenshots (right) ──
    #
    # LEFT — digital surface checklist (where the mark appears in the live app)
    add_text(s, "Digital surfaces",
             left=Inches(0.7), top=Inches(2.35), width=Inches(6), height=Inches(0.35),
             size=14, bold=True, font=DISPLAY, color=FG_HIGH)

    surfaces = [
        ("Favicon",        "vidiq_logo_black_bg.svg",   "Browser tab · bookmark · iOS home-screen"),
        ("Top nav",        "logo.tsx component",        "Persistent across every route · clickable home"),
        ("Hero splash",    "LogoSplash component",      "Animated session-gated landing animation"),
        ("OG / Twitter",   "vidiq_logo_white_bg.png",   "1200×630 social-share preview · all 4 routes"),
        ("JSON-LD",        "Organization schema",       "Logo URL embedded in structured data"),
        ("Slide footer",   "vidiq_logo_black_bg.png",   "0.32\" mark on every one of 33 slides"),
        ("FB Page cover",  "IMG-8.png · 1640×856",      "Cover photo + profile-pic mark on FB Page"),
    ]
    for i, (label, file, where) in enumerate(surfaces):
        y = Inches(2.78 + i * 0.55)
        add_card(s, left=Inches(0.7), top=y, width=Inches(6.0), height=Inches(0.48),
                 accent=ACCENT_VIOLET)
        # number badge
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(0.85), y + Inches(0.1),
                                    Inches(0.28), Inches(0.28))
        badge.fill.solid(); badge.fill.fore_color.rgb = ACCENT_VIOLET
        badge.line.fill.background()
        add_text(s, str(i + 1),
                 left=Inches(0.85), top=y + Inches(0.12),
                 width=Inches(0.28), height=Inches(0.25),
                 size=10, bold=True, font=DISPLAY, color=BG_DEEP, align=PP_ALIGN.CENTER)
        add_text(s, label,
                 left=Inches(1.25), top=y + Inches(0.06),
                 width=Inches(1.7), height=Inches(0.25),
                 size=10, bold=True, font=DISPLAY, color=FG_HIGH)
        add_text(s, file,
                 left=Inches(2.95), top=y + Inches(0.07),
                 width=Inches(2.0), height=Inches(0.25),
                 size=8, font=MONO, color=ACCENT_VIOLET)
        add_text(s, where,
                 left=Inches(1.25), top=y + Inches(0.27),
                 width=Inches(4.9), height=Inches(0.2),
                 size=8, font=BODY, color=FG_MID)

    # ── RIGHT — applied screenshots ─────────────────────────────────
    add_text(s, "Applied — screenshots",
             left=Inches(7.0), top=Inches(2.35), width=Inches(5.65), height=Inches(0.35),
             size=14, bold=True, font=DISPLAY, color=FG_HIGH)

    # FB Page cover (IMG-8) — top
    cover = asset("IMG-8.png")
    if cover:
        add_card(s, left=Inches(7.0), top=Inches(2.78),
                 width=Inches(5.65), height=Inches(2.0), accent=ACCENT_FUCHSIA)
        add_text(s, "FB PAGE COVER · 1640 × 856",
                 left=Inches(7.15), top=Inches(2.85),
                 width=Inches(5.4), height=Inches(0.22),
                 size=8, bold=True, font=BODY, color=ACCENT_FUCHSIA)
        add_image(s, cover,
                  left=Inches(7.15), top=Inches(3.1),
                  width=Inches(5.35), height=Inches(1.6),
                  border_color=ACCENT_FUCHSIA, border_pt=0.6)

    # Two side-by-side: FB page profile + IG welcome
    fb_page = asset("fb-page.png")
    ig_welcome = asset("insta-welcome-post.jpeg", "insta-welcome-post.png")
    bw = Inches(2.78); bh = Inches(1.95); bg = Inches(0.09)
    bsy = Inches(4.95)
    # FB page profile
    if fb_page:
        bx = Inches(7.0)
        add_card(s, left=bx, top=bsy, width=bw, height=bh, accent=ACCENT_VIOLET)
        add_text(s, "FB · fb.com/…id=61588939576479",
                 left=bx + Inches(0.12), top=bsy + Inches(0.07),
                 width=bw - Inches(0.24), height=Inches(0.22),
                 size=8, bold=True, font=BODY, color=ACCENT_VIOLET)
        add_image(s, fb_page,
                  left=bx + Inches(0.12), top=bsy + Inches(0.3),
                  width=bw - Inches(0.24), height=Inches(1.55),
                  border_color=ACCENT_VIOLET, border_pt=0.5)
    # IG welcome
    if ig_welcome:
        bx = Inches(7.0) + bw + bg
        add_card(s, left=bx, top=bsy, width=bw, height=bh, accent=ACCENT_CYAN)
        add_text(s, "IG · @vidiq_official",
                 left=bx + Inches(0.12), top=bsy + Inches(0.07),
                 width=bw - Inches(0.24), height=Inches(0.22),
                 size=8, bold=True, font=BODY, color=ACCENT_CYAN)
        add_image(s, ig_welcome,
                  left=bx + Inches(0.12), top=bsy + Inches(0.3),
                  width=bw - Inches(0.24), height=Inches(1.55),
                  border_color=ACCENT_CYAN, border_pt=0.5)

    # Footer rationale
    add_text(s, "Result: 7 / 7 surfaces audited green · KPI #4 (Brand consistency) scores 5 / 5 in the rubric tracker.",
             left=Inches(0.7), top=Inches(7.05), width=Inches(12), height=Inches(0.3),
             size=9, font=BODY, color=FG_LOW, align=PP_ALIGN.CENTER)

    set_notes(s, "Show that the mark appears everywhere — same SVG file, no drift. Left column lists "
                 "each surface with the file/component name so the examiner can verify in the repo. "
                 "Right column shows the social-side screenshots where brand legibility is most "
                 "scrutinised (Facebook + Instagram). KPI #4 (Brand consistency) is therefore 5 / 5. "
                 "Time: 50 s.")


# ── Build ──────────────────────────────────────────────────────────────────


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_1_cover(prs)
    slide_2_problem(prs)
    slide_3_product(prs)
    slide_4_pipeline(prs)
    slide_5_brand(prs)
    slide_6_ad(prs)
    slide_7_social(prs)
    slide_8_meta_ads(prs)
    slide_9_conversation(prs)
    slide_10_keywords(prs)
    slide_11_seo(prs)
    slide_12_google_ads(prs)
    slide_13_budget(prs)
    slide_14_competitive(prs)              # Section A + Section C
    slide_15_competitive_social(prs)       # Section B (FB + IG)
    slide_16_swot(prs)                     # SWOT 2x2 grid
    slide_17_kpi_tracker_full(prs)         # All 18 KPIs in detail
    slide_18_kpi_scorecard(prs)            # 6x3 visual summary
    slide_19_qa_backup(prs)                # Q&A backup

    # ── Digital-marketing appendix ────────────────────────────────────────
    slide_20_strategy_overview(prs)        # 5 pillars · positioning · GTM thesis
    slide_21_brand_identity(prs)           # Palette + typography + voice
    slide_22_video_ad_campaign(prs)        # 30-s ad · 6-beat script · 4 cutdowns
    slide_23_content_calendar(prs)         # 14-day · pillars · format mix
    slide_24_meta_ads_4_sets(prs)          # 4 ad-set deep dive with aim per set
    slide_25_conversation_layer(prs)       # Welcome · FAQs · saved replies
    slide_26_keyword_seo(prs)              # 18 KWs + on-page audit
    slide_27_google_ads(prs)               # 3 campaigns deep dive
    slide_28_budget(prs)                   # $740 distribution + phasing
    slide_29_measurement(prs)              # Pixel/CAPI · gtag · KPI targets
    slide_30_moats_roadmap(prs)            # 5 moats + 8-phase execution
    slide_31_seo_evidence(prs)             # GSC + sitemap + Lighthouse screenshots
    slide_32_logo_system(prs)              # 3 logo lockups · large · brand-rules strip
    slide_33_logo_in_the_wild(prs)         # Logo applied across 7 surfaces · screenshots

    out_dir = Path(__file__).resolve().parent / "submissions"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_file = out_dir / "VidIQ_Final_Presentation.pptx"
    try:
        prs.save(out_file)
    except PermissionError:
        # PowerPoint has the file open → write to a versioned fallback
        i = 2
        while (out_dir / f"VidIQ_Final_Presentation_v{i}.pptx").exists():
            i += 1
        out_file = out_dir / f"VidIQ_Final_Presentation_v{i}.pptx"
        prs.save(out_file)
        print(f"[warn] primary file locked - wrote {out_file.name} instead")
    print(f"[ok] wrote {out_file} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
